#!/usr/bin/env python3
"""
Tool: Fix PDSM Presentation
Descricao: Corrige claim do slide 7 e gera relatorio de densidade da apresentacao PDSM 4.0

Uso:
    py -3.13 tools/fix_pdsm_presentation.py [--input PATH] [--output PATH] [--dry-run]

Inputs:
    --input: Caminho do PPTX de entrada (default: exports/PDSM_Extrema_4.0_Apresentacao_2026.pptx)
    --output: Caminho do PPTX de saida (default: <input>_v2.pptx)
    --dry-run: Apenas relatar problemas sem corrigir

Outputs:
    - PPTX corrigido (se nao --dry-run)
    - Relatorio de correcoes e densidade (stdout)

Dependencias:
    - python-pptx
"""

import argparse
import sys
import io
from pathlib import Path

from pptx import Presentation

# Encoding fix para Windows console
if sys.platform == "win32":
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

PROJECT_ROOT = Path(__file__).parent.parent

# Substituicoes de texto: (slide_index_0based, texto_antigo, texto_novo)
TEXT_CORRECTIONS = [
    (
        6,
        "CLI-02 v2.7 avaliado como tecnicamente SUPERIOR às referências internacionais:",
        "CLI-02 v2.7 avaliado como tecnicamente alinhado e complementar às referências internacionais:",
    ),
]

DENSITY_THRESHOLD = 15


def find_and_replace_text(slide, old_text, new_text):
    """Substitui texto exato em shapes de um slide. Retorna True se encontrou."""
    replaced = False
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            if para.text == old_text:
                # Preservar formatacao: alterar texto do primeiro run, limpar os demais
                if para.runs:
                    para.runs[0].text = new_text
                    for run in para.runs[1:]:
                        run.text = ""
                else:
                    para.text = new_text
                replaced = True
    return replaced


def density_report(prs):
    """Gera relatorio de densidade por slide."""
    results = []
    for i, slide in enumerate(prs.slides):
        count = len(slide.shapes)
        severity = None
        if count > 25:
            severity = "ALTO"
        elif count > DENSITY_THRESHOLD:
            severity = "MEDIO"
        results.append({
            "slide": i + 1,
            "shapes": count,
            "severity": severity,
        })
    return results


def main():
    parser = argparse.ArgumentParser(
        description="Corrige claim e gera relatorio de densidade da apresentacao PDSM 4.0"
    )
    parser.add_argument(
        "--input",
        type=str,
        default=str(PROJECT_ROOT / "exports" / "PDSM_Extrema_4.0_Apresentacao_2026.pptx"),
        help="Caminho do PPTX de entrada",
    )
    parser.add_argument(
        "--output",
        type=str,
        default=None,
        help="Caminho do PPTX de saida (default: <input>_v2.pptx)",
    )
    parser.add_argument(
        "--dry-run",
        action="store_true",
        help="Apenas relatar problemas sem corrigir",
    )
    args = parser.parse_args()

    input_path = Path(args.input)
    if not input_path.exists():
        print(f"[ERRO] Arquivo nao encontrado: {input_path}")
        sys.exit(1)

    if args.output:
        output_path = Path(args.output)
    else:
        output_path = input_path.with_name(input_path.stem + "_v2.pptx")

    prs = Presentation(str(input_path))

    # --- Correcoes de texto ---
    print("=== CORRECOES DE TEXTO ===\n")
    corrections_applied = 0

    for slide_idx, old_text, new_text in TEXT_CORRECTIONS:
        slide = prs.slides[slide_idx]
        slide_num = slide_idx + 1

        # Verificar se o texto existe
        found = any(
            para.text == old_text
            for shape in slide.shapes
            if shape.has_text_frame
            for para in shape.text_frame.paragraphs
        )

        if found:
            if args.dry_run:
                print(f"  [DETECTADO] Slide {slide_num}:")
                print(f"    DE:   \"{old_text}\"")
                print(f"    PARA: \"{new_text}\"")
            else:
                replaced = find_and_replace_text(slide, old_text, new_text)
                if replaced:
                    print(f"  [CORRIGIDO] Slide {slide_num}:")
                    print(f"    DE:   \"{old_text}\"")
                    print(f"    PARA: \"{new_text}\"")
                    corrections_applied += 1
        else:
            print(f"  [NAO ENCONTRADO] Slide {slide_num}: \"{old_text[:60]}...\"")

    # --- Relatorio de densidade ---
    print("\n=== RELATORIO DE DENSIDADE ===\n")
    density = density_report(prs)
    high_density_count = 0

    for item in density:
        marker = ""
        if item["severity"] == "ALTO":
            marker = " << ALTO (considerar dividir)"
            high_density_count += 1
        elif item["severity"] == "MEDIO":
            marker = " < MEDIO"
            high_density_count += 1
        print(f"  Slide {item['slide']:2d}: {item['shapes']:2d} shapes{marker}")

    print(f"\n  Resumo: {high_density_count}/{len(density)} slides acima de {DENSITY_THRESHOLD} shapes")

    # --- Salvar ---
    if not args.dry_run and corrections_applied > 0:
        prs.save(str(output_path))
        print(f"\n=== SALVO: {output_path} ===")
        print(f"  {corrections_applied} correcao(oes) aplicada(s)")
    elif args.dry_run:
        print(f"\n=== DRY RUN: nenhuma alteracao salva ===")
    else:
        print(f"\n=== Nenhuma correcao necessaria ===")


if __name__ == "__main__":
    main()
