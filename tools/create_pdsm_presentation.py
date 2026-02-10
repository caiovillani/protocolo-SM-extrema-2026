# -*- coding: utf-8 -*-
"""
Gerador de Apresentação PDSM Extrema 4.0 - Etapa 2026
Formato: PowerPoint (.pptx)
Público: Gestores da Secretaria de Saúde de Extrema/MG
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

# Alias para compatibilidade com o código existente
RgbColor = RGBColor

# Cores institucionais
AZUL_SUS = RgbColor(0, 92, 169)       # #005CA9
VERDE_SAUDE = RgbColor(0, 168, 89)    # #00A859
CINZA_ESCURO = RgbColor(51, 51, 51)   # #333333
CINZA_CLARO = RgbColor(240, 240, 240) # #F0F0F0
BRANCO = RgbColor(255, 255, 255)
VERMELHO = RgbColor(220, 53, 69)      # Para alertas
LARANJA = RgbColor(255, 153, 0)       # Para prioridades
AMARELO = RgbColor(255, 193, 7)


def create_presentation():
    """Cria a apresentação completa."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # Widescreen 16:9
    prs.slide_height = Inches(7.5)

    # BLOCO 1: ABERTURA
    create_title_slide(prs)
    create_agenda_slide(prs)
    create_timeline_slide(prs)

    # BLOCO 2: RESGATE 2025
    create_diagnostico_slide(prs)
    create_conquistas_2025_slide(prs)
    create_linha_cuidado_tea_slide(prs)
    create_protocolos_desenvolvidos_slide(prs)
    create_metodologia_slide(prs)

    # BLOCO 3: VISÃO 2026 - METAS
    create_transicao_modelo_slide(prs)
    create_metas_estrategicas_slide(prs)
    create_classificacao_risco_slide(prs)
    create_indicadores_metas_slide(prs)

    # BLOCO 4: MACRO AÇÕES 1º SEMESTRE
    create_expansao_infraestrutura_slide(prs)
    create_compras_licitacoes_slide(prs)
    create_equipe_capsi_slide(prs)
    create_educacao_permanente_slide(prs)
    create_integracoes_slide(prs)
    create_governanca_slide(prs)

    # BLOCO 5: EXPANSÃO RAPS 2026-2028
    create_roadmap_raps_slide(prs)
    create_investimentos_slide(prs)
    create_rede_intersetorial_slide(prs)

    # BLOCO 6: PROTOCOLO DE COMPARTILHAMENTO
    create_modelo_integrado_slide(prs)
    create_fluxo_operacional_slide(prs)
    create_contrarreferencia_slide(prs)

    # BLOCO 7: ENCERRAMENTO
    create_cronograma_critico_slide(prs)
    create_proximos_passos_slide(prs)

    return prs


def add_slide_with_title(prs, title, subtitle=None):
    """Adiciona slide com título padronizado."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Barra superior azul
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.8))
    bar.fill.solid()
    bar.fill.fore_color.rgb = AZUL_SUS
    bar.line.fill.background()

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = BRANCO

    # Subtítulo (se houver)
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(0.3))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(14)
        p.font.color.rgb = RgbColor(200, 220, 255)

    # Linha verde inferior decorativa
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.8), prs.slide_width, Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = VERDE_SAUDE
    line.line.fill.background()

    return slide


def add_bullet_text(slide, left, top, width, height, items, font_size=16):
    """Adiciona lista com bullets."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = CINZA_ESCURO
        p.space_after = Pt(8)

    return box


def add_table(slide, left, top, width, height, headers, rows, col_widths=None):
    """Adiciona tabela formatada."""
    num_rows = len(rows) + 1
    num_cols = len(headers)

    table_shape = slide.shapes.add_table(num_rows, num_cols, Inches(left), Inches(top), Inches(width), Inches(height))
    table = table_shape.table

    # Definir larguras das colunas
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    # Cabeçalho
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = AZUL_SUS
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = BRANCO
        p.alignment = PP_ALIGN.CENTER

    # Dados
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(value)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.color.rgb = CINZA_ESCURO
            # Alternar cores de fundo
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CINZA_CLARO

    return table_shape


# ============================================================================
# BLOCO 1: ABERTURA
# ============================================================================

def create_title_slide(prs):
    """Slide 1: Capa."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Fundo gradiente (barra azul grande)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = AZUL_SUS
    bg.line.fill.background()

    # Linha verde decorativa
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.2), prs.slide_width, Inches(0.15))
    line.fill.solid()
    line.fill.fore_color.rgb = VERDE_SAUDE
    line.line.fill.background()

    # Título principal
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "PDSM Extrema 4.0"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = BRANCO
    p.alignment = PP_ALIGN.CENTER

    # Subtítulo
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Plano de Desenvolvimento de Saúde Mental"
    p.font.size = Pt(28)
    p.font.color.rgb = RgbColor(200, 220, 255)
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "Etapa 2026"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = VERDE_SAUDE
    p.alignment = PP_ALIGN.CENTER

    # Rodapé
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.5))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "Secretaria Municipal de Saúde | Extrema/MG | Janeiro 2026"
    p.font.size = Pt(14)
    p.font.color.rgb = RgbColor(180, 200, 230)
    p.alignment = PP_ALIGN.CENTER


def create_agenda_slide(prs):
    """Slide 2: Agenda."""
    slide = add_slide_with_title(prs, "Agenda", "Estrutura da Apresentação")

    items = [
        ("1", "Resgate 2025", "Diagnóstico, conquistas e fundamentos construídos"),
        ("2", "Visão 2026 - Metas Estratégicas", "Transição de modelo e indicadores"),
        ("3", "Macro Ações 1º Semestre", "Infraestrutura, equipes, capacitação"),
        ("4", "Expansão RAPS 2026-2028", "Roadmap, investimentos, rede intersetorial"),
        ("5", "Protocolo de Compartilhamento", "Modelo integrado APS-CSM"),
        ("6", "Cronograma e Próximos Passos", "Entregas críticas e aprovações"),
    ]

    y = 1.2
    for num, title, desc in items:
        # Número em círculo
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), Inches(y), Inches(0.5), Inches(0.5))
        circle.fill.solid()
        circle.fill.fore_color.rgb = VERDE_SAUDE
        circle.line.fill.background()

        num_box = slide.shapes.add_textbox(Inches(0.8), Inches(y + 0.05), Inches(0.5), Inches(0.4))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = BRANCO
        p.alignment = PP_ALIGN.CENTER

        # Título do item
        title_box = slide.shapes.add_textbox(Inches(1.5), Inches(y), Inches(5), Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = CINZA_ESCURO

        # Descrição
        desc_box = slide.shapes.add_textbox(Inches(1.5), Inches(y + 0.35), Inches(10), Inches(0.4))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = RgbColor(100, 100, 100)

        y += 0.9


def create_timeline_slide(prs):
    """Slide 3: Linha do Tempo 2025-2028."""
    slide = add_slide_with_title(prs, "Visão de Longo Prazo", "RAPS 1.0 → RAPS 4.0 (2025-2028)")

    # Linha horizontal
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.5), Inches(11), Inches(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = AZUL_SUS
    line.line.fill.background()

    phases = [
        ("2025", "Fase 1: Fundação", ["Diagnóstico situacional", "Educação Permanente", "Protocolos base", "Linha Cuidado TEA/DI"], CINZA_ESCURO),
        ("2026", "Fase 2: Expansão", ["CAPS I → CAPS II", "CECO II", "CAPS-i (infância)", "Centro Integrar 4.0"], VERDE_SAUDE),
        ("2027", "Fase 3: Especialização", ["CAPS AD II", "Consolidação matriciamento", "Integração RUE"], AZUL_SUS),
        ("2028", "Fase 4: Maturidade", ["UAA", "LSMHG (4-8 leitos)", "RAPS 4.0 completa"], RgbColor(128, 0, 128)),
    ]

    x = 1.2
    for year, title, items, color in phases:
        # Marcador
        marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 1), Inches(3.35), Inches(0.4), Inches(0.4))
        marker.fill.solid()
        marker.fill.fore_color.rgb = color
        marker.line.fill.background()

        # Ano
        year_box = slide.shapes.add_textbox(Inches(x), Inches(2.8), Inches(2.5), Inches(0.4))
        tf = year_box.text_frame
        p = tf.paragraphs[0]
        p.text = year
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

        # Título da fase
        title_box = slide.shapes.add_textbox(Inches(x), Inches(4), Inches(2.5), Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = CINZA_ESCURO
        p.alignment = PP_ALIGN.CENTER

        # Items
        items_box = slide.shapes.add_textbox(Inches(x), Inches(4.4), Inches(2.5), Inches(2.5))
        tf = items_box.text_frame
        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(10)
            p.font.color.rgb = CINZA_ESCURO
            p.alignment = PP_ALIGN.CENTER

        x += 3


# ============================================================================
# BLOCO 2: RESGATE 2025
# ============================================================================

def create_diagnostico_slide(prs):
    """Slide 4: Diagnóstico Situacional."""
    slide = add_slide_with_title(prs, "Diagnóstico Situacional", "RAPS 1.0 - Problemas Identificados (Set/2025)")

    problems = [
        ("CAPS I", ["Indicadores inconsistentes", "Baixa tecnologia de informação", "Fraca integração com rede"]),
        ("Núcleo Psicologia", ["Sobrecarga", "Casos mistos (leves a graves)", "Baixo matriciamento"]),
        ("Centro Integrar", ["Escopo indefinido", "Sem protocolos padronizados", "Isolado da RAPS"]),
        ("Emergência/Hospital", ["Baixa capacitação em SM", "Complicações iatrogênicas", "Desconectado da RAPS"]),
        ("eMulti APS", ["Escopo indefinido", "Sem matriciamento estruturado", "Baixa resolutividade"]),
    ]

    x = 0.5
    for title, items in problems:
        # Caixa do serviço
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.2), Inches(2.4), Inches(2.8))
        box.fill.solid()
        box.fill.fore_color.rgb = CINZA_CLARO
        box.line.color.rgb = VERMELHO
        box.line.width = Pt(2)

        # Título
        title_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(1.3), Inches(2.2), Inches(0.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = VERMELHO
        p.alignment = PP_ALIGN.CENTER

        # Problemas
        items_box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(1.8), Inches(2.1), Inches(2))
        tf = items_box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"✗ {item}"
            p.font.size = Pt(11)
            p.font.color.rgb = CINZA_ESCURO
            p.space_after = Pt(6)

        x += 2.55

    # Nota de rodapé
    note = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.4))
    tf = note.text_frame
    p = tf.paragraphs[0]
    p.text = "Fonte: Diagnóstico situacional PDSM Extrema 4.0 (Set/2025)"
    p.font.size = Pt(10)
    p.font.italic = True
    p.font.color.rgb = RgbColor(128, 128, 128)


def create_conquistas_2025_slide(prs):
    """Slide 5: Conquistas Set-Dez 2025."""
    slide = add_slide_with_title(prs, "Conquistas Set-Dez 2025", "Fundamentos Construídos")

    # Coluna esquerda
    left_items = [
        "Etapa 1 da Linha de Cuidado TEA/DI iniciada",
        "Educação Permanente estruturada em 8 eixos temáticos",
        "Meta: 40h/ano por profissional (~150 profissionais)",
        "Descompressão da Rede: acolhimentos e mutirão",
    ]
    add_bullet_text(slide, 0.5, 1.2, 5.5, 3, left_items, font_size=15)

    # Coluna direita
    right_items = [
        "Produtos de navegação RAS construídos",
        "Articulações institucionais estabelecidas",
        "Auditoria metodológica CI/PDSM (12 pontos corrigidos)",
        "Protocolo piloto de Compartilhamento do Cuidado",
    ]
    add_bullet_text(slide, 6.5, 1.2, 5.5, 3, right_items, font_size=15)

    # Destaque central
    highlight = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(4.3), Inches(9), Inches(1.5))
    highlight.fill.solid()
    highlight.fill.fore_color.rgb = RgbColor(230, 247, 230)
    highlight.line.color.rgb = VERDE_SAUDE
    highlight.line.width = Pt(2)

    highlight_text = slide.shapes.add_textbox(Inches(2.3), Inches(4.5), Inches(8.5), Inches(1.2))
    tf = highlight_text.text_frame
    p = tf.paragraphs[0]
    p.text = "Transição iniciada: Modelo Reativo → Modelo Proativo"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = VERDE_SAUDE
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "Vigilância ativa | Diagnóstico territorial | Estratificação de risco | Busca ativa"
    p.font.size = Pt(14)
    p.font.color.rgb = CINZA_ESCURO
    p.alignment = PP_ALIGN.CENTER


def create_linha_cuidado_tea_slide(prs):
    """Slide 6: Linha de Cuidado TEA/DI."""
    slide = add_slide_with_title(prs, "Linha de Cuidado TEA/DI", "Framework Consolidado (5 Eixos Estruturantes)")

    eixos = [
        ("I", "Intervenção Precoce", "Estimulação 0-3a (Precoce I), 3-6a (Precoce II)"),
        ("II", "Inclusão e Autonomia", "Formação professores, habilidades sociais, atividades adaptadas"),
        ("III", "Apoio Familiar", "Programas educativos, cuidado ao cuidador, destigmatização"),
        ("IV", "Gestão de Rede", "Metodologia de matriciamento, coordenação PTS"),
        ("V", "Acesso e Equidade", "Equipes móveis, telessaúde, transporte médico rural"),
    ]

    y = 1.1
    for num, title, desc in eixos:
        # Número
        num_box = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(y), Inches(0.6), Inches(0.6))
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = AZUL_SUS
        num_box.line.fill.background()

        num_text = slide.shapes.add_textbox(Inches(0.5), Inches(y + 0.1), Inches(0.6), Inches(0.4))
        tf = num_text.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = BRANCO
        p.alignment = PP_ALIGN.CENTER

        # Título
        title_box = slide.shapes.add_textbox(Inches(1.3), Inches(y), Inches(3.5), Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = CINZA_ESCURO

        # Descrição
        desc_box = slide.shapes.add_textbox(Inches(1.3), Inches(y + 0.4), Inches(5), Inches(0.5))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = RgbColor(100, 100, 100)

        y += 1

    # Sistema P1/P2/P3
    p_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(1.1), Inches(5.5), Inches(4.5))
    p_box.fill.solid()
    p_box.fill.fore_color.rgb = CINZA_CLARO
    p_box.line.color.rgb = AZUL_SUS

    p_title = slide.shapes.add_textbox(Inches(7.2), Inches(1.3), Inches(5), Inches(0.4))
    tf = p_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Sistema de Priorização"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = AZUL_SUS
    p.alignment = PP_ALIGN.CENTER

    headers = ["Prioridade", "Prazo", "Critérios"]
    rows = [
        ("P1 - Urgente", "30 dias", "<3a rastreio+ | Regressão | Risco"),
        ("P2 - Alta", "90 dias", "3-6a suspeita | Comprometimento"),
        ("P3 - Regular", "180 dias", ">6a suspeita | Adultos tardio"),
    ]
    add_table(slide, 7.2, 1.9, 5.1, 2.5, headers, rows, [1.3, 1.1, 2.7])


def create_protocolos_desenvolvidos_slide(prs):
    """Slide 7: Protocolos Desenvolvidos."""
    slide = add_slide_with_title(prs, "Documentação Técnica Desenvolvida", "7.143 linhas de protocolos normativos")

    headers = ["Código", "Protocolo", "Linhas", "Status"]
    rows = [
        ("CLI-02", "Transtorno do Espectro Autista (v2.7)", "1.928", "✓ Normativo"),
        ("PCC-06", "Macrofluxo Narrativo DI/TEA", "1.461", "✓ Normativo"),
        ("GN-01", "Guia Narrativo APS DI/TEA", "1.759", "✓ Normativo"),
        ("CLI-04", "Álcool e Outras Drogas", "941", "✓ Normativo"),
        ("ANA-01", "Análise Integrada TEA/DI", "1.054", "✓ Validado"),
    ]
    add_table(slide, 0.5, 1.2, 12, 3, headers, rows, [1.5, 5.5, 1.5, 2])

    # Destaque de validação
    valid_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.5), Inches(12), Inches(1.8))
    valid_box.fill.solid()
    valid_box.fill.fore_color.rgb = RgbColor(230, 245, 255)
    valid_box.line.color.rgb = AZUL_SUS

    valid_title = slide.shapes.add_textbox(Inches(0.7), Inches(4.7), Inches(11.5), Inches(0.4))
    tf = valid_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Validação Técnica Internacional"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = AZUL_SUS

    valid_text = slide.shapes.add_textbox(Inches(0.7), Inches(5.1), Inches(11.5), Inches(1))
    tf = valid_text.text_frame
    p = tf.paragraphs[0]
    p.text = "CLI-02 v2.7 avaliado como tecnicamente alinhado e complementar às referências internacionais:"
    p.font.size = Pt(13)
    p.font.color.rgb = CINZA_ESCURO

    p = tf.add_paragraph()
    p.text = "• NICE CG128 (UK)  •  Einstein Care Pathway  •  CDC Milestones (EUA)"
    p.font.size = Pt(13)
    p.font.color.rgb = VERDE_SAUDE
    p.font.bold = True


def create_metodologia_slide(prs):
    """Slide 8: Metodologia Aplicada."""
    slide = add_slide_with_title(prs, "Metodologia Aplicada", "Planifica SUS + Ciência da Implementação")

    # Planifica SUS
    planifica_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(5.8), Inches(3.5))
    planifica_box.fill.solid()
    planifica_box.fill.fore_color.rgb = RgbColor(240, 248, 255)
    planifica_box.line.color.rgb = AZUL_SUS

    planifica_title = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(5.4), Inches(0.4))
    tf = planifica_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Planifica SUS (CONASS/MS)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = AZUL_SUS

    planifica_items = [
        "4 Macro Etapas estruturadas",
        "12+ Micro Etapas detalhadas",
        "Formato narrativo (O que | Por que | Como)",
        "Matriz RACI de responsabilidades",
    ]
    add_bullet_text(slide, 0.7, 1.9, 5.4, 2.5, planifica_items, font_size=13)

    # CFIR
    cfir_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(1.2), Inches(5.8), Inches(3.5))
    cfir_box.fill.solid()
    cfir_box.fill.fore_color.rgb = RgbColor(245, 255, 240)
    cfir_box.line.color.rgb = VERDE_SAUDE

    cfir_title = slide.shapes.add_textbox(Inches(6.7), Inches(1.4), Inches(5.4), Inches(0.4))
    tf = cfir_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Ciência da Implementação (CFIR)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = VERDE_SAUDE

    cfir_items = [
        "5 domínios de implementação",
        "Auditoria com 12 correções críticas",
        "Monitoramento de consequências",
        "Modelo lógico participativo",
    ]
    add_bullet_text(slide, 6.7, 1.9, 5.4, 2.5, cfir_items, font_size=13)

    # Instrumentos
    instr_title = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(12), Inches(0.4))
    tf = instr_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Instrumentos Validados Implementados"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = CINZA_ESCURO

    instr_text = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12), Inches(0.8))
    tf = instr_text.text_frame
    p = tf.paragraphs[0]
    p.text = "M-CHAT-R/F (85-95% sens.)  |  IRDI (31 indicadores)  |  CuidaSM  |  IFBrM  |  PHQ-9  |  GAD-7  |  Columbia-SSRS"
    p.font.size = Pt(14)
    p.font.color.rgb = CINZA_ESCURO
    p.alignment = PP_ALIGN.CENTER


# ============================================================================
# BLOCO 3: VISÃO 2026 - METAS
# ============================================================================

def create_transicao_modelo_slide(prs):
    """Slide 9: Transição de Modelo."""
    slide = add_slide_with_title(prs, "Transição de Modelo", "Reativo → Proativo")

    # Modelo Reativo (antes)
    before_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(5.5), Inches(3))
    before_box.fill.solid()
    before_box.fill.fore_color.rgb = RgbColor(255, 240, 240)
    before_box.line.color.rgb = VERMELHO
    before_box.line.width = Pt(2)

    before_title = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(5.1), Inches(0.4))
    tf = before_title.text_frame
    p = tf.paragraphs[0]
    p.text = "ANTES: Modelo Reativo"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = VERMELHO

    before_items = [
        "Espera pela demanda espontânea",
        "Atendimento por crise",
        "Fragmentação do cuidado",
        "Baixa integração de rede",
    ]
    add_bullet_text(slide, 0.7, 1.9, 5.1, 2, before_items, font_size=14)

    # Seta
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.1), Inches(2.4), Inches(0.8), Inches(0.6))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = VERDE_SAUDE
    arrow.line.fill.background()

    # Modelo Proativo (depois)
    after_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(1.2), Inches(5.5), Inches(3))
    after_box.fill.solid()
    after_box.fill.fore_color.rgb = RgbColor(240, 255, 240)
    after_box.line.color.rgb = VERDE_SAUDE
    after_box.line.width = Pt(2)

    after_title = slide.shapes.add_textbox(Inches(7.2), Inches(1.4), Inches(5.1), Inches(0.4))
    tf = after_title.text_frame
    p = tf.paragraphs[0]
    p.text = "AGORA: Modelo Proativo"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = VERDE_SAUDE

    after_items = [
        "Vigilância ativa via ACS",
        "Diagnóstico territorial sistemático",
        "Estratificação de risco (PHQ-9, GAD-7)",
        "Busca ativa de casos",
    ]
    add_bullet_text(slide, 7.2, 1.9, 5.1, 2, after_items, font_size=14)

    # Resultados esperados
    result_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.5), Inches(12), Inches(2))
    result_box.fill.solid()
    result_box.fill.fore_color.rgb = RgbColor(245, 250, 255)
    result_box.line.color.rgb = AZUL_SUS

    result_title = slide.shapes.add_textbox(Inches(0.7), Inches(4.7), Inches(11.5), Inches(0.4))
    tf = result_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Resultados Esperados"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = AZUL_SUS

    result_text = slide.shapes.add_textbox(Inches(0.7), Inches(5.1), Inches(11.5), Inches(1.2))
    tf = result_text.text_frame
    p = tf.paragraphs[0]
    p.text = "↓ Crises evitadas  |  ↓ Internações  |  ↑ Diagnóstico precoce  |  ↑ Resolutividade APS  |  ↑ Qualidade de vida"
    p.font.size = Pt(15)
    p.font.color.rgb = CINZA_ESCURO
    p.alignment = PP_ALIGN.CENTER


def create_metas_estrategicas_slide(prs):
    """Slide 10: 10 Metas Estratégicas 2026."""
    slide = add_slide_with_title(prs, "10 Metas Estratégicas 2026", "Direcionadores do Plano")

    metas = [
        "Modernização das Tecnologias Municipais de Cuidado",
        "Fortalecimento da RCPD (Rede de Cuidados à Pessoa com Deficiência)",
        "Modernização e Ampliação da RAPS",
        "Aproximação com SES-MG",
        "Amadurecimento das Relações Intersetoriais",
        "Ampliação das Ações em Rede",
        "Fortalecimento das e-Multi/APS e e-ESF/APS",
        "Apoio à estruturação de Políticas de Cuidado no Território",
        "Educação Permanente direcionada a Governança",
        "Fortalecimento do Modelo Psicossocial de Atenção",
    ]

    # Duas colunas
    for i, meta in enumerate(metas[:5]):
        y = 1.1 + (i * 0.95)
        num_box = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(y), Inches(0.45), Inches(0.45))
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = VERDE_SAUDE
        num_box.line.fill.background()

        num_text = slide.shapes.add_textbox(Inches(0.5), Inches(y + 0.07), Inches(0.45), Inches(0.35))
        tf = num_text.text_frame
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = BRANCO
        p.alignment = PP_ALIGN.CENTER

        meta_box = slide.shapes.add_textbox(Inches(1.1), Inches(y), Inches(5), Inches(0.8))
        tf = meta_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = meta
        p.font.size = Pt(13)
        p.font.color.rgb = CINZA_ESCURO

    for i, meta in enumerate(metas[5:]):
        y = 1.1 + (i * 0.95)
        num_box = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.5), Inches(y), Inches(0.45), Inches(0.45))
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = AZUL_SUS
        num_box.line.fill.background()

        num_text = slide.shapes.add_textbox(Inches(6.5), Inches(y + 0.07), Inches(0.45), Inches(0.35))
        tf = num_text.text_frame
        p = tf.paragraphs[0]
        p.text = str(i + 6)
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = BRANCO
        p.alignment = PP_ALIGN.CENTER

        meta_box = slide.shapes.add_textbox(Inches(7.1), Inches(y), Inches(5.5), Inches(0.8))
        tf = meta_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = meta
        p.font.size = Pt(13)
        p.font.color.rgb = CINZA_ESCURO


def create_classificacao_risco_slide(prs):
    """Slide 11: Sistema de Classificação de Risco."""
    slide = add_slide_with_title(prs, "Sistema de Classificação de Risco", "Matriz de Cores para Priorização")

    # Tabela de cores
    colors_data = [
        ("VERMELHO", "Crise/Emergência", "PS + CAPS + APS", VERMELHO),
        ("LARANJA", "Alto", "CAPS II + APS", LARANJA),
        ("AMARELO", "Moderado", "CSM + APS", AMARELO),
        ("VERDE", "Baixo", "APS + matriciamento", VERDE_SAUDE),
        ("AZUL", "Crônico estável", "APS longitudinal", AZUL_SUS),
    ]

    y = 1.2
    for cor, risco, local, rgb_color in colors_data:
        # Indicador de cor
        color_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(1.5), Inches(0.7))
        color_box.fill.solid()
        color_box.fill.fore_color.rgb = rgb_color
        color_box.line.fill.background()

        color_text = slide.shapes.add_textbox(Inches(0.5), Inches(y + 0.15), Inches(1.5), Inches(0.4))
        tf = color_text.text_frame
        p = tf.paragraphs[0]
        p.text = cor
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = BRANCO if cor not in ["AMARELO"] else CINZA_ESCURO
        p.alignment = PP_ALIGN.CENTER

        # Risco
        risco_box = slide.shapes.add_textbox(Inches(2.2), Inches(y + 0.15), Inches(2.5), Inches(0.4))
        tf = risco_box.text_frame
        p = tf.paragraphs[0]
        p.text = risco
        p.font.size = Pt(14)
        p.font.color.rgb = CINZA_ESCURO

        # Local
        local_box = slide.shapes.add_textbox(Inches(4.8), Inches(y + 0.15), Inches(3), Inches(0.4))
        tf = local_box.text_frame
        p = tf.paragraphs[0]
        p.text = local
        p.font.size = Pt(14)
        p.font.color.rgb = CINZA_ESCURO

        y += 0.85

    # Critérios VERMELHO
    crit_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8), Inches(1.2), Inches(4.5), Inches(4))
    crit_box.fill.solid()
    crit_box.fill.fore_color.rgb = RgbColor(255, 245, 245)
    crit_box.line.color.rgb = VERMELHO

    crit_title = slide.shapes.add_textbox(Inches(8.2), Inches(1.4), Inches(4.1), Inches(0.4))
    tf = crit_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Critérios VERMELHO (Emergência)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = VERMELHO

    crit_items = [
        "Risco suicida agudo (plano + meios + intenção)",
        "Risco de agressão (auto ou hetero)",
        "Sintomas psicóticos agudos com agitação",
        "Risco de exposição moral",
        "→ Ação: SAMU 192 imediato",
    ]
    add_bullet_text(slide, 8.2, 1.9, 4.1, 3, crit_items, font_size=11)


def create_indicadores_metas_slide(prs):
    """Slide 12: Indicadores e Metas 2026."""
    slide = add_slide_with_title(prs, "Indicadores e Metas 2026", "Monitoramento Contínuo")

    headers = ["Indicador", "Meta", "Periodicidade"]
    rows = [
        ("Taxa de classificação de risco", "100%", "Mensal"),
        ("PTS elaborado em ≤90 dias", "100%", "Trimestral"),
        ("Tempo médio de espera", "≤3 dias", "Mensal"),
        ("Reavaliação PTS no prazo", "≥90%", "Trimestral"),
        ("Contrarreferência efetiva", "100%", "Mensal"),
        ("Taxa de abandono", "≤15%", "Trimestral"),
        ("Cobertura M-CHAT (18-24m)", "≥80%", "Mensal"),
        ("Diagnóstico TEA P1 em ≤30d", "≥90%", "Trimestral"),
    ]
    add_table(slide, 0.5, 1.2, 12, 4.5, headers, rows, [6, 2.5, 2.5])


# ============================================================================
# BLOCO 4: MACRO AÇÕES 1º SEMESTRE 2026
# ============================================================================

def create_expansao_infraestrutura_slide(prs):
    """Slide 13: Expansão da Infraestrutura."""
    slide = add_slide_with_title(prs, "Expansão da Infraestrutura", "1º Semestre 2026")

    infra_items = [
        ("CAPS-I → CAPS-II", "Habilitação e ampliação de serviços", AZUL_SUS),
        ("Reforma CAPS", "Adequação física e acessibilidade", AZUL_SUS),
        ("Ampliação Centro Integrar", "Expansão de capacidade (TEA/DI)", VERDE_SAUDE),
        ("Inauguração CAPS-i", "Infância e juventude - novo equipamento", VERDE_SAUDE),
    ]

    y = 1.2
    for title, desc, color in infra_items:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(12), Inches(1))
        box.fill.solid()
        box.fill.fore_color.rgb = CINZA_CLARO
        box.line.color.rgb = color
        box.line.width = Pt(2)

        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(y + 0.15), Inches(4), Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = color

        desc_box = slide.shapes.add_textbox(Inches(5), Inches(y + 0.25), Inches(7), Inches(0.5))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = CINZA_ESCURO

        y += 1.2

    # Nota sobre CECO II
    note_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.3), Inches(12), Inches(1.2))
    note_box.fill.solid()
    note_box.fill.fore_color.rgb = RgbColor(255, 250, 230)
    note_box.line.color.rgb = LARANJA

    note_title = slide.shapes.add_textbox(Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.4))
    tf = note_title.text_frame
    p = tf.paragraphs[0]
    p.text = "NOVO: CECO II (Portaria 5.738/2024)"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = LARANJA

    note_text = slide.shapes.add_textbox(Inches(0.8), Inches(5.9), Inches(11.5), Inches(0.5))
    tf = note_text.text_frame
    p = tf.paragraphs[0]
    p.text = "Dispositivo federal: baixo investimento, alto impacto. Custo de implantação coberto por portaria."
    p.font.size = Pt(12)
    p.font.color.rgb = CINZA_ESCURO


def create_compras_licitacoes_slide(prs):
    """Slide 14: Compras e Licitações."""
    slide = add_slide_with_title(prs, "Compras e Licitações", "Investimentos Planejados - Coordenações de Saúde Mental")

    items_left = [
        "Móveis para equipamentos",
        "Eletroeletrônicos",
        "Materiais de construção",
        "Sala Sensorial (Centro Integrar)",
        "Testes Psicométricos (ADOS-2, ADI-R, CARS-2)",
    ]

    items_right = [
        "TVs para atendimento",
        "Materiais de Teleatendimento (parceria APS)",
        "Veículos (CAPS-i)",
        "Materiais de apoio (jogos educativos)",
        "Equipamentos administrativos",
    ]

    add_bullet_text(slide, 0.5, 1.2, 5.5, 4, items_left, font_size=15)
    add_bullet_text(slide, 6.5, 1.2, 5.5, 4, items_right, font_size=15)

    # Nota de ação
    action_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.2), Inches(12), Inches(1.2))
    action_box.fill.solid()
    action_box.fill.fore_color.rgb = RgbColor(255, 245, 235)
    action_box.line.color.rgb = LARANJA

    action_text = slide.shapes.add_textbox(Inches(0.8), Inches(5.4), Inches(11.5), Inches(0.8))
    tf = action_text.text_frame
    p = tf.paragraphs[0]
    p.text = "AÇÃO NECESSÁRIA: Aprovação orçamentária para processos licitatórios"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = LARANJA
    p.alignment = PP_ALIGN.CENTER


def create_equipe_capsi_slide(prs):
    """Slide 15: Equipe CAPS-i."""
    slide = add_slide_with_title(prs, "Equipe CAPS-i", "Cronograma de Implantação - Infância e Juventude")

    # Timeline
    phases = [
        ("Seleção", "Jan-Fev", "Processo seletivo equipe técnica"),
        ("Contratação", "Fev-Mar", "Formalização de vínculos"),
        ("Treinamento", "Mar-Abr", "Capacitação em neurodesenvolvimento"),
        ("Inauguração", "Abr-Mai", "Início das atividades"),
    ]

    x = 0.8
    for phase, period, desc in phases:
        # Caixa da fase
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.5), Inches(2.8), Inches(2.5))
        box.fill.solid()
        box.fill.fore_color.rgb = CINZA_CLARO
        box.line.color.rgb = AZUL_SUS

        # Título
        title_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(1.7), Inches(2.6), Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = phase
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = AZUL_SUS
        p.alignment = PP_ALIGN.CENTER

        # Período
        period_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(2.1), Inches(2.6), Inches(0.4))
        tf = period_box.text_frame
        p = tf.paragraphs[0]
        p.text = period
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = VERDE_SAUDE
        p.alignment = PP_ALIGN.CENTER

        # Descrição
        desc_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(2.6), Inches(2.6), Inches(1.2))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = CINZA_ESCURO
        p.alignment = PP_ALIGN.CENTER

        x += 3.1

    # Custeio federal
    custeio_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.5), Inches(12), Inches(2))
    custeio_box.fill.solid()
    custeio_box.fill.fore_color.rgb = RgbColor(235, 255, 235)
    custeio_box.line.color.rgb = VERDE_SAUDE

    custeio_title = slide.shapes.add_textbox(Inches(0.8), Inches(4.7), Inches(11.5), Inches(0.4))
    tf = custeio_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Custeio Federal Mensal (Portaria 5.500/2024)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = VERDE_SAUDE

    custeio_text = slide.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(11.5), Inches(1))
    tf = custeio_text.text_frame
    p = tf.paragraphs[0]
    p.text = "CAPS-i (Infantojuvenil): R$ 48.797/mês"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = CINZA_ESCURO
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "Primeiro reajuste em 12 anos (+19,5%)"
    p.font.size = Pt(14)
    p.font.color.rgb = RgbColor(100, 100, 100)
    p.alignment = PP_ALIGN.CENTER


def create_educacao_permanente_slide(prs):
    """Slide 16: Educação Permanente."""
    slide = add_slide_with_title(prs, "Educação Permanente Especializada", "Neurodesenvolvimento e Governança")

    # Dr. Rubens
    dr_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(5.8), Inches(2.5))
    dr_box.fill.solid()
    dr_box.fill.fore_color.rgb = RgbColor(240, 248, 255)
    dr_box.line.color.rgb = AZUL_SUS

    dr_title = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(5.4), Inches(0.4))
    tf = dr_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Contratação: Dr. Rubens W."
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = AZUL_SUS

    dr_items = [
        "Especialista em Neurodesenvolvimento",
        "Capacitação das equipes RAPS",
        "Supervisão clínica em TEA/DI",
        "Formação de multiplicadores",
    ]
    add_bullet_text(slide, 0.7, 1.9, 5.4, 2, dr_items, font_size=13)

    # Governança
    gov_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(1.2), Inches(5.8), Inches(2.5))
    gov_box.fill.solid()
    gov_box.fill.fore_color.rgb = RgbColor(245, 255, 240)
    gov_box.line.color.rgb = VERDE_SAUDE

    gov_title = slide.shapes.add_textbox(Inches(6.7), Inches(1.4), Inches(5.4), Inches(0.4))
    tf = gov_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Governança para Coordenações"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = VERDE_SAUDE

    gov_items = [
        "Capacitação em planejamento",
        "Estratégias de gestão local",
        "Indicadores e monitoramento",
        "Conselhos Gestores Locais",
    ]
    add_bullet_text(slide, 6.7, 1.9, 5.4, 2, gov_items, font_size=13)

    # Programa geral
    prog_title = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12), Inches(0.4))
    tf = prog_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Programa Estruturado (8 Eixos Temáticos)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = CINZA_ESCURO

    eixos = [
        "Depressão", "Ansiedade/TOC", "Psicoses", "Uso de substâncias",
        "Neurodesenvolvimento", "Infância/adolescência", "Crise/suicídio", "Intersetorialidade"
    ]

    eixos_text = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12), Inches(0.5))
    tf = eixos_text.text_frame
    p = tf.paragraphs[0]
    p.text = "  |  ".join(eixos)
    p.font.size = Pt(12)
    p.font.color.rgb = CINZA_ESCURO
    p.alignment = PP_ALIGN.CENTER

    # Meta
    meta_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3), Inches(5.2), Inches(7), Inches(1.2))
    meta_box.fill.solid()
    meta_box.fill.fore_color.rgb = VERDE_SAUDE
    meta_box.line.fill.background()

    meta_text = slide.shapes.add_textbox(Inches(3.2), Inches(5.4), Inches(6.6), Inches(0.8))
    tf = meta_text.text_frame
    p = tf.paragraphs[0]
    p.text = "Meta: 40h/ano por profissional"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = BRANCO
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "~150 profissionais (RAPS + APS regional)"
    p.font.size = Pt(14)
    p.font.color.rgb = RgbColor(220, 255, 220)
    p.alignment = PP_ALIGN.CENTER


def create_integracoes_slide(prs):
    """Slide 17: Integrações Estratégicas."""
    slide = add_slide_with_title(prs, "Integrações Estratégicas", "Rede de Atenção e Eventos")

    # Integrações
    integs = [
        ("Matriciamento RAPS", "Encontros sistemáticos APS ↔ CAPS ↔ CSM"),
        ("Integração RUE/RAPS", "Urgência/Emergência com fluxos SM"),
        ("Hospital São Lucas/RAPS", "Capacitação e protocolos integrados"),
        ("Teleconsulta 0800 644 6543", "Suporte federal para casos complexos"),
    ]

    y = 1.2
    for title, desc in integs:
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(4), Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"→ {title}"
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = AZUL_SUS

        desc_box = slide.shapes.add_textbox(Inches(0.8), Inches(y + 0.35), Inches(5.5), Inches(0.4))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = CINZA_ESCURO

        y += 0.85

    # Eventos
    events_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(1.2), Inches(5.8), Inches(3))
    events_box.fill.solid()
    events_box.fill.fore_color.rgb = RgbColor(255, 250, 240)
    events_box.line.color.rgb = LARANJA

    events_title = slide.shapes.add_textbox(Inches(6.7), Inches(1.4), Inches(5.4), Inches(0.4))
    tf = events_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Eventos 1º Semestre"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = LARANJA

    events_items = [
        "ABRIL: Mês do TEA/DI",
        "   • Conscientização e inclusão",
        "   • Ações comunitárias",
        "MAIO: Luta Antimanicomial",
        "   • Destigmatização",
        "   • Direitos em saúde mental",
    ]

    y = 1.9
    for item in events_items:
        item_box = slide.shapes.add_textbox(Inches(6.9), Inches(y), Inches(5), Inches(0.35))
        tf = item_box.text_frame
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(12) if item.startswith("   ") else Pt(14)
        p.font.bold = not item.startswith("   ")
        p.font.color.rgb = CINZA_ESCURO
        y += 0.33


def create_governanca_slide(prs):
    """Slide 18: Governança Descentralizada."""
    slide = add_slide_with_title(prs, "Governança Descentralizada", "Estrutura de Gestão Participativa")

    # Conselhos Gestores
    cg_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(6), Inches(3.5))
    cg_box.fill.solid()
    cg_box.fill.fore_color.rgb = RgbColor(240, 248, 255)
    cg_box.line.color.rgb = AZUL_SUS

    cg_title = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(5.6), Inches(0.4))
    tf = cg_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Conselhos Gestores Locais"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = AZUL_SUS

    cg_items = [
        "Meta: 100% dos equipamentos até 01/06/2026",
        "Composição (Lei 8.142/1990):",
        "   • 10% Gestão",
        "   • 40% Trabalhadores",
        "   • 50% Usuários/Familiares",
        "Deliberações sobre serviço local",
    ]

    y = 1.9
    for item in cg_items:
        item_box = slide.shapes.add_textbox(Inches(0.7), Inches(y), Inches(5.6), Inches(0.35))
        tf = item_box.text_frame
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(12) if item.startswith("   ") else Pt(13)
        p.font.color.rgb = CINZA_ESCURO
        y += 0.35

    # Reuniões
    reunioes_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.7), Inches(1.2), Inches(5.8), Inches(3.5))
    reunioes_box.fill.solid()
    reunioes_box.fill.fore_color.rgb = RgbColor(245, 255, 240)
    reunioes_box.line.color.rgb = VERDE_SAUDE

    reunioes_title = slide.shapes.add_textbox(Inches(6.9), Inches(1.4), Inches(5.4), Inches(0.4))
    tf = reunioes_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Estrutura de Reuniões"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = VERDE_SAUDE

    reunioes_items = [
        "Semanal: Coordenação Municipal",
        "   • Segunda-feira, 16h",
        "   • Casos, fluxos, EP",
        "Trimestral: Departamento (100%)",
        "   • Indicadores e metas",
        "   • Planejamento estratégico",
    ]

    y = 1.9
    for item in reunioes_items:
        item_box = slide.shapes.add_textbox(Inches(6.9), Inches(y), Inches(5.4), Inches(0.35))
        tf = item_box.text_frame
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(12) if item.startswith("   ") else Pt(13)
        p.font.color.rgb = CINZA_ESCURO
        y += 0.35


# ============================================================================
# BLOCO 5: EXPANSÃO RAPS 2026-2028
# ============================================================================

def create_roadmap_raps_slide(prs):
    """Slide 19: Roadmap RAPS 1.0 → 4.0."""
    slide = add_slide_with_title(prs, "Roadmap RAPS 1.0 → 4.0", "Expansão Planejada 2025-2028")

    headers = ["Fase", "Ano", "Dispositivos", "Investimento"]
    rows = [
        ("Fase 1", "2025", "Fundação (diagnóstico, EP, protocolos)", "—"),
        ("Fase 2", "2026", "CAPS II, CECO II, CAPS-i", "R$ 30.000 + federal"),
        ("Fase 3", "2027", "CAPS AD II", "R$ 50.000"),
        ("Fase 4", "2028", "UAA, LSMHG (4-8 leitos)", "R$ 103.000"),
    ]
    add_table(slide, 0.5, 1.2, 12, 2.5, headers, rows, [1.5, 1.5, 5.5, 3])

    # Visão RAPS 4.0
    vision_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4), Inches(12), Inches(2.5))
    vision_box.fill.solid()
    vision_box.fill.fore_color.rgb = RgbColor(235, 245, 255)
    vision_box.line.color.rgb = AZUL_SUS

    vision_title = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(11.5), Inches(0.4))
    tf = vision_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Visão RAPS 4.0 (2028)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = AZUL_SUS

    vision_items = [
        "CAPS II (qualificado) + CAPS-i (infância) + CAPS AD II (álcool/drogas)",
        "CECO II + UAA (Unidade de Acolhimento Adulto)",
        "Centro de Saúde Mental + Centro Integrar 4.0 (neurodivergências)",
        "LSMHG (4-8 leitos em hospital geral) + Emergência capacitada",
        "eMulti APS matriciadora + Rede intersetorial articulada",
    ]
    add_bullet_text(slide, 0.8, 4.7, 11.5, 2, vision_items, font_size=12)


def create_investimentos_slide(prs):
    """Slide 20: Investimentos e Sustentabilidade."""
    slide = add_slide_with_title(prs, "Investimentos e Sustentabilidade Financeira", "Portarias 5.500/2024 e 5.738/2024")

    # Aumento federal
    aumento_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.2), Inches(5.8), Inches(1.5))
    aumento_box.fill.solid()
    aumento_box.fill.fore_color.rgb = VERDE_SAUDE
    aumento_box.line.fill.background()

    aumento_text = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(5.4), Inches(1.2))
    tf = aumento_text.text_frame
    p = tf.paragraphs[0]
    p.text = "+19,5%"
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = BRANCO
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "Primeiro reajuste em 12 anos"
    p.font.size = Pt(14)
    p.font.color.rgb = RgbColor(220, 255, 220)
    p.alignment = PP_ALIGN.CENTER

    # Tabela de valores
    headers = ["CAPS", "Anterior", "Atualizado"]
    rows = [
        ("CAPS I → II", "R$ 33.086", "R$ 50.257"),
        ("CAPS-i (Infantojuvenil)", "R$ 32.130", "R$ 48.797"),
        ("CAPS AD", "R$ 39.780", "R$ 60.388"),
    ]
    add_table(slide, 6.5, 1.2, 5.8, 2.2, headers, rows, [2.5, 1.5, 1.5])

    # Investimento total
    total_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(3), Inches(12), Inches(1.5))
    total_box.fill.solid()
    total_box.fill.fore_color.rgb = RgbColor(245, 250, 255)
    total_box.line.color.rgb = AZUL_SUS

    total_title = slide.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(11.5), Inches(0.4))
    tf = total_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Investimento Total Estimado (2025-2028)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = AZUL_SUS

    total_text = slide.shapes.add_textbox(Inches(0.8), Inches(3.7), Inches(11.5), Inches(0.6))
    tf = total_text.text_frame
    p = tf.paragraphs[0]
    p.text = "CAPEX: ~R$ 213.000  |  OPEX Federal Adicional: ~R$ 500.000+/ano até 2028"
    p.font.size = Pt(18)
    p.font.color.rgb = CINZA_ESCURO
    p.alignment = PP_ALIGN.CENTER

    # Detalhamento
    detail_title = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(12), Inches(0.4))
    tf = detail_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Detalhamento por Fase"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = CINZA_ESCURO

    detail_items = [
        "2026: CAPS I→II (R$ 30.000) + CECO II (federal)",
        "2027: CAPS AD II (R$ 50.000) + CAPS-i (R$ 30.000)",
        "2028: UAA (R$ 70.000) + LSMHG (R$ 33.000)",
    ]
    add_bullet_text(slide, 0.5, 5.2, 12, 1.5, detail_items, font_size=12)


def create_rede_intersetorial_slide(prs):
    """Slide 21: Rede Intersetorial Formalizada."""
    slide = add_slide_with_title(prs, "Rede Intersetorial Formalizada", "Articulação Ampliada")

    setores = [
        ("Educação", "PSE, formação professores, PEI inclusivo"),
        ("CRAS/CREAS", "Proteção social, BPC, vulnerabilidade"),
        ("Conselho Tutelar", "Direitos, notificações, violência"),
        ("Esportes/Cultura", "Inclusão comunitária, atividades adaptadas"),
        ("3º Setor", "Associações, grupos de apoio, economia solidária"),
        ("Justiça/MP", "Internação involuntária, direito à saúde"),
    ]

    x = 0.5
    y = 1.2
    for i, (setor, desc) in enumerate(setores):
        if i == 3:
            x = 6.5
            y = 1.2

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(5.8), Inches(1))
        box.fill.solid()
        box.fill.fore_color.rgb = CINZA_CLARO
        box.line.color.rgb = AZUL_SUS if i < 3 else VERDE_SAUDE

        setor_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.1), Inches(5.4), Inches(0.4))
        tf = setor_box.text_frame
        p = tf.paragraphs[0]
        p.text = setor
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = AZUL_SUS if i < 3 else VERDE_SAUDE

        desc_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.5), Inches(5.4), Inches(0.4))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = CINZA_ESCURO

        y += 1.2

    # Comitê Gestor
    comite_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.8), Inches(12), Inches(1.6))
    comite_box.fill.solid()
    comite_box.fill.fore_color.rgb = RgbColor(255, 250, 235)
    comite_box.line.color.rgb = LARANJA

    comite_title = slide.shapes.add_textbox(Inches(0.8), Inches(5), Inches(11.5), Inches(0.4))
    tf = comite_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Comitê Gestor Intersetorial (permanente)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = LARANJA

    comite_text = slide.shapes.add_textbox(Inches(0.8), Inches(5.4), Inches(11.5), Inches(0.8))
    tf = comite_text.text_frame
    p = tf.paragraphs[0]
    p.text = "Representantes: Saúde Mental + Educação + CRAS/CREAS + Habitação + Desenvolvimento + Direitos Humanos + Esporte"
    p.font.size = Pt(12)
    p.font.color.rgb = CINZA_ESCURO
    p.alignment = PP_ALIGN.CENTER


# ============================================================================
# BLOCO 6: PROTOCOLO DE COMPARTILHAMENTO
# ============================================================================

def create_modelo_integrado_slide(prs):
    """Slide 22: Modelo APS-CSM Integrado."""
    slide = add_slide_with_title(prs, "Modelo APS-CSM Integrado", "Princípios do Compartilhamento do Cuidado")

    principios = [
        ("Cuidado Integral", "Abordagem biopsicossocial completa"),
        ("Equidade", "Acesso universal com priorização por necessidade"),
        ("Compartilhamento", "Não é encaminhamento, é cuidado conjunto"),
        ("Coordenação APS", "A APS mantém a responsabilidade longitudinal"),
        ("Matriciamento", "Regulador de qualidade da rede"),
        ("PTS Compartilhado", "Ferramenta central de gestão do cuidado"),
    ]

    x = 0.5
    y = 1.2
    for i, (title, desc) in enumerate(principios):
        if i == 3:
            x = 6.5
            y = 1.2

        num_box = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.5), Inches(0.5))
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = VERDE_SAUDE
        num_box.line.fill.background()

        num_text = slide.shapes.add_textbox(Inches(x), Inches(y + 0.08), Inches(0.5), Inches(0.35))
        tf = num_text.text_frame
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = BRANCO
        p.alignment = PP_ALIGN.CENTER

        title_box = slide.shapes.add_textbox(Inches(x + 0.6), Inches(y), Inches(5), Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = CINZA_ESCURO

        desc_box = slide.shapes.add_textbox(Inches(x + 0.6), Inches(y + 0.4), Inches(5), Inches(0.4))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = RgbColor(100, 100, 100)

        y += 1.1


def create_fluxo_operacional_slide(prs):
    """Slide 23: Fluxo Operacional em 7 Passos."""
    slide = add_slide_with_title(prs, "Fluxo Operacional em 7 Passos", "Jornada do Paciente no Compartilhamento")

    passos = [
        ("1", "Entrada e Acolhimento", "Qualquer porta de entrada"),
        ("2", "Avaliação e-ESF", "Exame mental, PHQ-9/GAD-7, classificação"),
        ("3", "Revisão e-Multi", "Validação, teleconsulta 0800"),
        ("4", "Matriciamento Intra-APS", "Discussão de equipe, PTS momentos 2-3"),
        ("5", "Solicitação NIRSM-R", "Se necessário, formulário padronizado"),
        ("6", "Cuidado Compartilhado", "CSM/CAPS, PTS continuidade"),
        ("7", "Alta/Reclassificação", "Estabilização ≥3 meses, retorno APS"),
    ]

    y = 1.1
    for num, title, desc in passos:
        # Número
        num_box = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(y), Inches(0.45), Inches(0.45))
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = AZUL_SUS
        num_box.line.fill.background()

        num_text = slide.shapes.add_textbox(Inches(0.5), Inches(y + 0.07), Inches(0.45), Inches(0.35))
        tf = num_text.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = BRANCO
        p.alignment = PP_ALIGN.CENTER

        # Título
        title_box = slide.shapes.add_textbox(Inches(1.1), Inches(y), Inches(4), Inches(0.4))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = CINZA_ESCURO

        # Descrição
        desc_box = slide.shapes.add_textbox(Inches(5.2), Inches(y + 0.05), Inches(7), Inches(0.35))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = RgbColor(100, 100, 100)

        y += 0.7


def create_contrarreferencia_slide(prs):
    """Slide 24: Critérios de Contrarreferência."""
    slide = add_slide_with_title(prs, "Critérios de Contrarreferência", "Retorno Estruturado à APS")

    headers = ["Critério", "Responsável", "Verificação"]
    rows = [
        ("Estabilização sustentada ≥3 meses", "CSM/CAPS", "Prontuário + avaliação clínica"),
        ("Metas do PTS alcançadas", "Equipe referência", "Checklist PTS"),
        ("Classificação verde/azul atingida", "Médico responsável", "Reavaliação formal"),
        ("Plano de seguimento APS definido", "APS + AE", "Documento compartilhado"),
        ("Família orientada sobre sinais de alerta", "e-Multi", "Registro + assinatura"),
    ]
    add_table(slide, 0.5, 1.2, 12, 3.5, headers, rows, [5, 2.5, 4])

    # Nota importante
    note_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5), Inches(12), Inches(1.5))
    note_box.fill.solid()
    note_box.fill.fore_color.rgb = RgbColor(245, 255, 240)
    note_box.line.color.rgb = VERDE_SAUDE

    note_title = slide.shapes.add_textbox(Inches(0.8), Inches(5.2), Inches(11.5), Inches(0.4))
    tf = note_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Princípio Fundamental"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = VERDE_SAUDE

    note_text = slide.shapes.add_textbox(Inches(0.8), Inches(5.6), Inches(11.5), Inches(0.7))
    tf = note_text.text_frame
    p = tf.paragraphs[0]
    p.text = "A contrarreferência não é 'dar alta', é transferir a coordenação do cuidado de volta para a APS com plano estruturado."
    p.font.size = Pt(13)
    p.font.color.rgb = CINZA_ESCURO
    p.alignment = PP_ALIGN.CENTER


# ============================================================================
# BLOCO 7: ENCERRAMENTO
# ============================================================================

def create_cronograma_critico_slide(prs):
    """Slide 25: Cronograma Crítico 1º Trimestre."""
    slide = add_slide_with_title(prs, "Cronograma Crítico", "Entregas do 1º Semestre 2026")

    headers = ["Data", "Entrega", "Status"]
    rows = [
        ("09/02", "Calendário de Educação Permanente", "Pendente"),
        ("15/02", "Baseline PTS (levantamento)", "Pendente"),
        ("01/03", "Início Regula Raps", "Pendente"),
        ("01/03", "PTS obrigatório para todos usuários", "Pendente"),
        ("15/03", "Planejamento campanhas (Abril/Maio)", "Pendente"),
        ("Abril", "Campanha TEA", "Planejado"),
        ("Maio", "Luta Antimanicomial", "Planejado"),
        ("01/06", "Conselhos Gestores Locais (100%)", "Pendente"),
        ("01/06", "Capacitação em planejamento (100%)", "Pendente"),
    ]
    add_table(slide, 0.5, 1.2, 12, 4.5, headers, rows, [1.5, 7, 2.5])


def create_proximos_passos_slide(prs):
    """Slide 26: Próximos Passos e Aprovação."""
    slide = add_slide_with_title(prs, "Próximos Passos", "Ações para Aprovação e Implementação")

    acoes = [
        ("Validação do cronograma de entregas", "Secretária + Coordenações"),
        ("Aprovação orçamentária (compras/licitações)", "Diretora Administrativa/Finanças"),
        ("Autorização para contratações (CAPS-i, Dr. Rubens)", "Secretária"),
        ("Definição de responsáveis por entrega (RACI)", "Coordenações"),
        ("Agendamento de reuniões de monitoramento", "Coordenação Municipal SM"),
    ]

    y = 1.2
    for acao, responsavel in acoes:
        # Checkbox
        check = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(0.4), Inches(0.4))
        check.fill.solid()
        check.fill.fore_color.rgb = BRANCO
        check.line.color.rgb = AZUL_SUS
        check.line.width = Pt(2)

        # Ação
        acao_box = slide.shapes.add_textbox(Inches(1.1), Inches(y), Inches(7), Inches(0.4))
        tf = acao_box.text_frame
        p = tf.paragraphs[0]
        p.text = acao
        p.font.size = Pt(14)
        p.font.color.rgb = CINZA_ESCURO

        # Responsável
        resp_box = slide.shapes.add_textbox(Inches(8.5), Inches(y), Inches(4), Inches(0.4))
        tf = resp_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"→ {responsavel}"
        p.font.size = Pt(12)
        p.font.color.rgb = AZUL_SUS
        p.font.italic = True

        y += 0.8

    # Call to action
    cta_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(5.5), Inches(9), Inches(1.3))
    cta_box.fill.solid()
    cta_box.fill.fore_color.rgb = VERDE_SAUDE
    cta_box.line.fill.background()

    cta_text = slide.shapes.add_textbox(Inches(2.2), Inches(5.7), Inches(8.6), Inches(1))
    tf = cta_text.text_frame
    p = tf.paragraphs[0]
    p.text = "Aprovação Solicitada"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = BRANCO
    p.alignment = PP_ALIGN.CENTER

    p = tf.add_paragraph()
    p.text = "PDSM Extrema 4.0 — Etapa 2026"
    p.font.size = Pt(16)
    p.font.color.rgb = RgbColor(220, 255, 220)
    p.alignment = PP_ALIGN.CENTER


def main():
    """Função principal."""
    print("Criando apresentação PDSM Extrema 4.0 - Etapa 2026...")

    prs = create_presentation()

    # Criar diretório de saída se não existir
    output_dir = Path(__file__).parent.parent / "exports"
    output_dir.mkdir(exist_ok=True)

    output_path = output_dir / "PDSM_Extrema_4.0_Apresentacao_2026.pptx"
    prs.save(str(output_path))

    print(f"Apresentação salva em: {output_path}")
    print(f"Total de slides: {len(prs.slides)}")

    return str(output_path)


if __name__ == "__main__":
    main()
