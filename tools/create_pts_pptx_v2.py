#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Gerador de PPTX para AULA-01: Projeto Terapêutico Singular v2.0

Este script cria uma apresentação PowerPoint com 36 slides baseada na
versão HTML da aula sobre PTS para médicos de família.

Requisitos: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Alias for consistency
RgbColor = RGBColor

# Cores do tema SUS
AZUL_SUS = RgbColor(0x00, 0x5C, 0xA9)
AZUL_ESCURO = RgbColor(0x00, 0x33, 0x66)
VERDE_SAUDE = RgbColor(0x00, 0xA8, 0x59)
VERDE_ESCURO = RgbColor(0x00, 0x7A, 0x42)
LARANJA_ALERTA = RgbColor(0xF5, 0x7C, 0x00)
VERMELHO_RISCO = RgbColor(0xD3, 0x2F, 0x2F)
ROXO_MOMENTO4 = RgbColor(0x7B, 0x1F, 0xA2)
CINZA_ESCURO = RgbColor(0x2D, 0x34, 0x36)
CINZA_MEDIO = RgbColor(0x63, 0x6E, 0x72)
BRANCO = RgbColor(0xFF, 0xFF, 0xFF)

# Cores de classificação de risco
RISCO_VERMELHO = RgbColor(0xD3, 0x2F, 0x2F)
RISCO_LARANJA = RgbColor(0xF5, 0x7C, 0x00)
RISCO_AMARELO = RgbColor(0xFF, 0xC1, 0x07)
RISCO_VERDE = RgbColor(0x4C, 0xAF, 0x50)
RISCO_AZUL = RgbColor(0x21, 0x96, 0xF3)


def set_shape_fill(shape, color):
    """Define cor de preenchimento de uma forma."""
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def add_title_slide(prs, title, subtitle=""):
    """Adiciona slide de título."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background azul
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    set_shape_fill(background, AZUL_SUS)
    background.line.fill.background()

    # Título
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(9), Inches(1.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = BRANCO
    p.alignment = PP_ALIGN.CENTER

    # Subtítulo
    if subtitle:
        sub_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4), Inches(9), Inches(1)
        )
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = BRANCO
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title, content_items, title_color=AZUL_SUS, note=""):
    """Adiciona slide com título e lista de itens."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Título
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = title_color
    p.alignment = PP_ALIGN.CENTER

    # Conteúdo
    content_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.3), Inches(9), Inches(5.5)
    )
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(item, tuple):
            # Item com formatação especial (texto, tamanho, negrito, cor)
            p.text = item[0]
            p.font.size = Pt(item[1] if len(item) > 1 else 20)
            p.font.bold = item[2] if len(item) > 2 else False
            if len(item) > 3 and item[3]:
                p.font.color.rgb = item[3]
        else:
            p.text = f"• {item}"
            p.font.size = Pt(20)

        p.font.color.rgb = CINZA_ESCURO
        p.space_after = Pt(12)

    # Nota do apresentador
    if note:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = note

    return slide


def add_case_slide(prs, name, age, details, scores="", classification="", classification_color=RISCO_AMARELO, note=""):
    """Adiciona slide de caso clínico."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Título
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Caso Clínico"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = AZUL_SUS
    p.alignment = PP_ALIGN.CENTER

    # Box do caso
    case_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1), Inches(9), Inches(4.5)
    )
    case_box.fill.solid()
    case_box.fill.fore_color.rgb = RgbColor(0xF5, 0xF6, 0xF7)
    case_box.line.color.rgb = LARANJA_ALERTA
    case_box.line.width = Pt(4)

    # Nome e idade
    name_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.2), Inches(8.6), Inches(0.5)
    )
    tf = name_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{name}, {age} anos"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = AZUL_ESCURO

    # Detalhes
    details_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.8), Inches(8.6), Inches(2.5)
    )
    tf = details_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = details
    p.font.size = Pt(18)
    p.font.color.rgb = CINZA_ESCURO
    p.line_spacing = 1.5

    # Escores
    if scores:
        scores_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(4.3), Inches(8.6), Inches(0.7)
        )
        scores_box.fill.solid()
        scores_box.fill.fore_color.rgb = RgbColor(0xE3, 0xF2, 0xFD)
        scores_box.line.fill.background()

        scores_text = slide.shapes.add_textbox(
            Inches(0.8), Inches(4.4), Inches(8.4), Inches(0.5)
        )
        tf = scores_text.text_frame
        p = tf.paragraphs[0]
        p.text = f"Escores: {scores}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = AZUL_SUS

    # Classificação
    if classification:
        class_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3), Inches(5.2), Inches(4), Inches(0.5)
        )
        set_shape_fill(class_box, classification_color)
        class_box.line.fill.background()

        class_text = slide.shapes.add_textbox(
            Inches(3), Inches(5.25), Inches(4), Inches(0.4)
        )
        tf = class_text.text_frame
        p = tf.paragraphs[0]
        p.text = classification
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = BRANCO if classification_color != RISCO_AMARELO else CINZA_ESCURO
        p.alignment = PP_ALIGN.CENTER

    if note:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = note

    return slide


def add_table_slide(prs, title, headers, rows, title_color=AZUL_SUS, note=""):
    """Adiciona slide com tabela."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Título
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = title_color
    p.alignment = PP_ALIGN.CENTER

    # Tabela
    cols = len(headers)
    table_rows = len(rows) + 1

    table = slide.shapes.add_table(
        table_rows, cols, Inches(0.3), Inches(1.1), Inches(9.4), Inches(0.5 * table_rows)
    ).table

    # Headers
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = AZUL_SUS
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = BRANCO

    # Rows
    for row_idx, row in enumerate(rows):
        for col_idx, cell_text in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(13)
            p.font.color.rgb = CINZA_ESCURO
            if row_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RgbColor(0xF5, 0xF6, 0xF7)

    if note:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = note

    return slide


def add_checkpoint_slide(prs, title, scenario, options, answer, note=""):
    """Adiciona slide de checkpoint interativo."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Título
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = LARANJA_ALERTA
    p.alignment = PP_ALIGN.CENTER

    # Cenário
    scenario_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1), Inches(9), Inches(1.5)
    )
    scenario_box.fill.solid()
    scenario_box.fill.fore_color.rgb = BRANCO
    scenario_box.line.color.rgb = LARANJA_ALERTA
    scenario_box.line.width = Pt(2)

    scenario_text = slide.shapes.add_textbox(
        Inches(0.6), Inches(1.1), Inches(8.8), Inches(1.3)
    )
    tf = scenario_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = scenario
    p.font.size = Pt(16)
    p.font.color.rgb = CINZA_ESCURO

    # Opções
    options_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.7), Inches(9), Inches(1.5)
    )
    tf = options_box.text_frame
    tf.word_wrap = True
    for i, opt in enumerate(options):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = opt
        p.font.size = Pt(16)
        p.font.color.rgb = CINZA_ESCURO
        p.space_after = Pt(8)

    # Resposta
    answer_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.5), Inches(9), Inches(1.2)
    )
    set_shape_fill(answer_box, VERDE_SAUDE)
    answer_box.line.fill.background()

    answer_text = slide.shapes.add_textbox(
        Inches(0.6), Inches(4.6), Inches(8.8), Inches(1)
    )
    tf = answer_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"Resposta: {answer}"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BRANCO

    if note:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = note

    return slide


def add_highlight_slide(prs, text, subtitle="", bg_color=VERDE_SAUDE, note=""):
    """Adiciona slide de destaque."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    set_shape_fill(bg, bg_color)
    bg.line.fill.background()

    # Texto principal
    text_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2), Inches(9), Inches(2)
    )
    tf = text_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = BRANCO
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        sub_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4), Inches(9), Inches(1)
        )
        tf = sub_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(20)
        p.font.color.rgb = BRANCO
        p.alignment = PP_ALIGN.CENTER

    if note:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = note

    return slide


def add_moments_slide(prs, note=""):
    """Adiciona slide dos 4 momentos."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Título
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Os 4 Momentos do PTS"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = AZUL_SUS
    p.alignment = PP_ALIGN.CENTER

    moments = [
        ("1", "Diagnóstico Integral", "Bio + Psico + Social", AZUL_SUS),
        ("2", "Definição de Metas", "SMART — com o usuário", VERDE_SAUDE),
        ("3", "Responsabilidades", "Quem faz o quê", LARANJA_ALERTA),
        ("4", "Reavaliação", "Documento vivo", ROXO_MOMENTO4),
    ]

    positions = [
        (Inches(0.5), Inches(1.5)),
        (Inches(5), Inches(1.5)),
        (Inches(0.5), Inches(3.8)),
        (Inches(5), Inches(3.8)),
    ]

    for i, (num, title, desc, color) in enumerate(moments):
        x, y = positions[i]

        # Card
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(4.3), Inches(2)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = BRANCO
        card.line.color.rgb = color
        card.line.width = Pt(3)

        # Número
        num_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x + Inches(1.7), y + Inches(0.2), Inches(0.8), Inches(0.8)
        )
        set_shape_fill(num_circle, color)
        num_circle.line.fill.background()

        num_text = slide.shapes.add_textbox(
            x + Inches(1.7), y + Inches(0.3), Inches(0.8), Inches(0.6)
        )
        tf = num_text.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = BRANCO
        p.alignment = PP_ALIGN.CENTER

        # Título do momento
        moment_title = slide.shapes.add_textbox(
            x + Inches(0.2), y + Inches(1.1), Inches(3.9), Inches(0.5)
        )
        tf = moment_title.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

        # Descrição
        moment_desc = slide.shapes.add_textbox(
            x + Inches(0.2), y + Inches(1.5), Inches(3.9), Inches(0.4)
        )
        tf = moment_desc.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = CINZA_MEDIO
        p.alignment = PP_ALIGN.CENTER

    if note:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = note

    return slide


def add_quiz_slide(prs, questions, note=""):
    """Adiciona slide de quiz."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Título
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.6)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Quiz Rápido — 4 Perguntas"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = AZUL_SUS
    p.alignment = PP_ALIGN.CENTER

    y_pos = Inches(1.1)
    for q, a in questions:
        # Pergunta
        q_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), y_pos, Inches(9.4), Inches(1.1)
        )
        q_box.fill.solid()
        q_box.fill.fore_color.rgb = RgbColor(0xF5, 0xF6, 0xF7)
        q_box.line.fill.background()

        q_text = slide.shapes.add_textbox(
            Inches(0.5), y_pos + Inches(0.1), Inches(9), Inches(0.5)
        )
        tf = q_text.text_frame
        p = tf.paragraphs[0]
        p.text = q
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = AZUL_ESCURO

        a_text = slide.shapes.add_textbox(
            Inches(0.5), y_pos + Inches(0.55), Inches(9), Inches(0.4)
        )
        tf = a_text.text_frame
        p = tf.paragraphs[0]
        p.text = f"→ {a}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = VERDE_SAUDE

        y_pos += Inches(1.2)

    if note:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = note

    return slide


def create_presentation():
    """Cria a apresentação completa."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ===== SLIDE 1: CAPA =====
    slide = add_title_slide(
        prs,
        "Projeto Terapêutico Singular",
        "Uma ferramenta para organizar o cuidado na APS"
    )
    # Footer
    footer = slide.shapes.add_textbox(Inches(2), Inches(5.8), Inches(6), Inches(0.8))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "Coordenação Municipal de Saúde Mental\nExtrema/MG — Janeiro 2026 | v2.0"
    p.font.size = Pt(14)
    p.font.color.rgb = BRANCO
    p.alignment = PP_ALIGN.CENTER

    # ===== SLIDE 2: CASO DONA MARIA =====
    add_case_slide(
        prs,
        "Dona Maria", "58",
        "DM2 + HAS há 15 anos\n\"Depressão\" há 3 anos\n8 consultas no último ano\nQueixa: \"dor no corpo\", \"não consigo dormir\"\nExames: normais\nSem resolução",
        "PHQ-9: 14 (moderada) | CuidaSM: Nível 3",
        "AMARELO - Moderado",
        RISCO_AMARELO,
        "Ler pausadamente. Enfatizar '8 consultas' e 'sem resolução'. PHQ-9 de 14 = depressão moderada. CuidaSM Nível 3 = PTS obrigatório."
    )

    # ===== SLIDE 3: PERGUNTA MOBILIZADORA =====
    add_highlight_slide(
        prs,
        "Quantas \"Donas Marias\"\nvocês têm na agenda?",
        "Casos que consomem múltiplas consultas sem resolução aparente",
        VERDE_ESCURO,
        "Deixar 30 segundos de silêncio. Aguardar respostas. Validar: 'Todo mundo tem uma Dona Maria'."
    )

    # ===== SLIDE 4: TRANSIÇÃO =====
    add_highlight_slide(
        prs,
        "Casos assim geram frustração\n— pra gente e pro paciente.",
        "O que vamos ver hoje é uma ferramenta que não resolve tudo,\nmas estrutura o cuidado.",
        VERDE_SAUDE,
        "Tom empático. Reconhecer a realidade da APS."
    )

    # ===== SLIDE 5: O QUE É PTS =====
    add_content_slide(
        prs,
        "O que é PTS?",
        [
            "Ferramenta de organização do cuidado",
            "Considera a singularidade de cada caso",
            "Construído COM o usuário, não PARA o usuário",
            "Atribui responsabilidades claras",
        ],
        AZUL_SUS,
        "Enfatizar: não é mais um papel para preencher. É uma lógica de organização."
    )

    # ===== SLIDE 6: PTS vs PLANO (PARTE 1) =====
    add_table_slide(
        prs,
        "PTS ≠ Plano de Cuidados Genérico",
        ["Plano Genérico", "PTS"],
        [
            ["Padronizado por condição", "Singularizado por pessoa"],
            ["Feito pela equipe", "Co-construído com usuário"],
            ["Foco na doença", "Foco na pessoa"],
        ],
        AZUL_SUS,
        "O PTS complementa para casos complexos."
    )

    # ===== SLIDE 7: PTS vs PLANO (PARTE 2) =====
    add_table_slide(
        prs,
        "PTS ≠ Plano de Cuidados Genérico",
        ["Plano Genérico", "PTS"],
        [
            ["Responsabilidades difusas", "Matriz clara de quem faz o quê"],
            ["Revisão eventual", "Reavaliação programada"],
            ["Usuário passivo", "Usuário corresponsável"],
        ],
        AZUL_SUS,
        "Destaque: 'todo mundo cuida = ninguém cuida'"
    )

    # ===== SLIDE 8: MARCOS LEGAIS =====
    add_content_slide(
        prs,
        "De onde vem o PTS?",
        [
            "Política Nacional de Humanização (PNH, 2004)",
            "Conceito de Clínica Ampliada e Compartilhada",
            "Cadernos de Atenção Básica nº 34 — Saúde Mental (MS, 2013)",
            "Linha de Cuidado TEA — MS Brasil 2025",
            "Princípio: cuidado centrado na pessoa, não na doença",
        ],
        AZUL_SUS,
        "Mostrar que é diretriz oficial, não modismo."
    )

    # ===== SLIDE 9: QUANDO ELABORAR + CuidaSM =====
    add_table_slide(
        prs,
        "Quando elaborar PTS? + CuidaSM",
        ["Nível CuidaSM", "Pontuação", "PTS"],
        [
            ["Nível 1", "0-1", "Opcional"],
            ["Nível 2", "2-3", "Recomendado"],
            ["Nível 3", "4-6", "OBRIGATÓRIO"],
            ["Nível 4", "7-11", "OBRIGATÓRIO + CAPS"],
        ],
        VERDE_SAUDE,
        "Dona Maria: CuidaSM Nível 3 = PTS obrigatório."
    )

    # ===== SLIDE 10: PTS PRELIMINAR =====
    add_content_slide(
        prs,
        "PTS Preliminar — Inovação MS 2025",
        [
            "Elaborado ANTES da confirmação diagnóstica",
            "Especialmente para TEA/DI na infância",
            "Base científica: janela de neuroplasticidade 0-3 anos",
            "Cada mês de espera é uma oportunidade perdida",
            ("\"A intervenção deve iniciar-se ANTES do diagnóstico formal\" — MS 2025", 16, False, CINZA_MEDIO),
        ],
        LARANJA_ALERTA,
        "Isso é mudança de paradigma. Vamos ver o caso Lucas."
    )

    # ===== SLIDE 11: PTS SIMPLIFICADO vs COMPLETO =====
    add_table_slide(
        prs,
        "PTS Simplificado vs Completo",
        ["Aspecto", "Simplificado", "Completo"],
        [
            ["Indicação", "Verde/Azul, TMC leve", "Amarelo+, TMG"],
            ["Diagnóstico", "Breve (1 página)", "Biopsicossocial"],
            ["Metas", "1-2 prioritárias", "3-5 com prazos"],
            ["Reavaliação", "30-60 dias", "Conforme complexidade"],
        ],
        AZUL_SUS,
        "Nem todo PTS precisa ser extenso."
    )

    # ===== SLIDE 12: OS 4 MOMENTOS =====
    add_moments_slide(
        prs,
        "São 4 momentos que formam um ciclo. As cores ajudam a identificar cada momento."
    )

    # ===== SLIDE 13: MOMENTO 1 CONCEITO =====
    add_content_slide(
        prs,
        "① Diagnóstico Integral",
        [
            "Vai além do CID-10",
            "Mapeia: vulnerabilidades + potencialidades",
            "Avaliação biopsicossocial completa",
            "Inclui: barreiras de acesso, rede de apoio",
            ("Não é diagnóstico médico — é diagnóstico SITUACIONAL", 18, True, VERDE_SAUDE),
        ],
        AZUL_SUS,
        "MFC já faz isso intuitivamente. PTS sistematiza."
    )

    # ===== SLIDE 14: CASO LUCAS =====
    add_case_slide(
        prs,
        "Lucas", "8",
        "Suspeita de TEA + sintomas TDAH\nM-CHAT-R/F positivo aos 24 meses (5/20)\nPrimeiras palavras: 30 meses (atraso)\nEscola: rendimento OK, isolamento social\nHipersensibilidade auditiva\nInteresse restrito: trens",
        "IRDI: 4 indicadores ausentes | CARS-2: Pendente",
        "LARANJA - P2 (90 dias) - PTS PRELIMINAR",
        RISCO_LARANJA,
        "Lucas ilustra PTS Preliminar. Intervenção não espera diagnóstico."
    )

    # ===== SLIDE 15: DIAGNÓSTICO DONA MARIA =====
    add_table_slide(
        prs,
        "Diagnóstico Integral — Dona Maria",
        ["Dimensão", "Avaliação"],
        [
            ["Biológica", "DM2 descompensada, HAS, insônia, PHQ-9: 14"],
            ["Psicológica", "Humor deprimido, luto não elaborado, baixa autoestima"],
            ["Social", "Mora sozinha, filhos distantes, renda insuficiente"],
            ["Potencialidades", "Vínculo com ACS, artesanato, boa adesão"],
        ],
        AZUL_SUS,
        "Não é 'paciente difícil' — é pessoa com múltiplas vulnerabilidades."
    )

    # ===== SLIDE 16: CHECKPOINT 1 =====
    add_checkpoint_slide(
        prs,
        "🤔 O que você faria?",
        "Paciente 45 anos, DM2 controlada, cansaço e desânimo há 2 meses. Primeira consulta. PHQ-9: 8 (leve)",
        [
            "A) Elaborar PTS completo imediatamente",
            "B) Solicitar exames e reavaliar em 2 semanas",
            "C) Encaminhar ao CAPS",
            "D) Aplicar CuidaSM e decidir baseado no nível",
        ],
        "D — PHQ-9 de 8 é depressão leve. CuidaSM determinará necessidade.",
        "Pausa de 2 minutos para discussão."
    )

    # ===== SLIDE 17: MOMENTO 2 CONCEITO =====
    add_content_slide(
        prs,
        "② Definição de Metas",
        [
            "Metas SMART: Específicas, Mensuráveis, Alcançáveis, Relevantes, Temporais",
            "Negociadas COM o usuário — não impostas",
            "Horizonte: 30d (curto) / 90d (médio) / >90d (longo)",
            "Foco em funcionalidade e qualidade de vida",
            ("❌ \"Melhorar a depressão\" → ✓ \"PHQ-9 < 10 em 90 dias\"", 16, False, CINZA_MEDIO),
        ],
        VERDE_SAUDE,
        "Meta não é 'controlar diabetes'. É algo que faça sentido para ela."
    )

    # ===== SLIDE 18: CASO CARLA =====
    add_case_slide(
        prs,
        "Carla", "32",
        "Depressão Maior (F32.2) + Ansiedade (F41.1)\n3 meses de acompanhamento compartilhado APS-CAPS\nMedicação: Sertralina 100mg\nPsicologia e-Multi: 8 sessões",
        "PHQ-9: 18→9 | GAD-7: 14→6 | CuidaSM: 3→2",
        "AMARELO → trajetória para VERDE",
        RISCO_AMARELO,
        "Carla ilustra compartilhamento APS-CAPS e evolução com escores."
    )

    # ===== SLIDE 19: METAS DONA MARIA =====
    add_table_slide(
        prs,
        "Metas — Dona Maria",
        ["Prazo", "Meta", "Indicador"],
        [
            ["30 dias", "Retomar sono de 6h/noite", "Autorrelato"],
            ["30 dias", "Café da manhã com vizinha", "Confirmação ACS"],
            ["90 dias", "Retornar ao artesanato", "Frequência ≥2x/mês"],
            ["90 dias", "PHQ-9 < 10", "Aplicação em consulta"],
        ],
        VERDE_SAUDE,
        "Metas sociais junto com metas clínicas. Negociadas com ela."
    )

    # ===== SLIDE 20: CHECKPOINT 2 =====
    add_checkpoint_slide(
        prs,
        "✏️ Transforme em SMART",
        "Meta original: \"Melhorar o sono da Dona Maria\"\n\nEm duplas (1 minuto): Reformulem usando critérios SMART.",
        [],
        "\"Retomar sono de 6h contínuas, verificado por autorrelato e vizinha, em 30 dias\"\nS: 6h contínuas | M: autorrelato + vizinha | A: realista | R: central | T: 30 dias",
        "Dar 1 minuto real para as duplas. Colher 2-3 respostas."
    )

    # ===== SLIDE 21: MOMENTO 3 CONCEITO =====
    add_content_slide(
        prs,
        "③ Divisão de Responsabilidades",
        [
            "QUEM faz O QUÊ, QUANDO",
            "Inclui: equipe + usuário + família + rede",
            "Define profissional de referência (gestor do caso)",
            ("Evita: \"todo mundo cuida = ninguém cuida\"", 20, True, LARANJA_ALERTA),
        ],
        LARANJA_ALERTA,
        "Este é o coração do PTS. Sem isso, vira boa intenção."
    )

    # ===== SLIDE 22: RESPONSABILIDADES EQUIPE =====
    add_table_slide(
        prs,
        "Responsabilidades — Equipe",
        ["Ator", "Ação", "Frequência"],
        [
            ["Médico(a)", "Revisão medicamentosa + PHQ-9", "Mensal"],
            ["Enfermeiro(a)", "Monitoramento DM/HAS + escuta", "Quinzenal"],
            ["ACS", "Visita + monitorar sono", "Semanal"],
            ["Psicólogo", "Acompanhamento luto (6 sessões)", "Quinzenal"],
        ],
        LARANJA_ALERTA,
        "Profissional de Referência: ACS Maria (melhor vínculo)."
    )

    # ===== SLIDE 23: RESPONSABILIDADES USUÁRIO/REDE =====
    add_table_slide(
        prs,
        "Responsabilidades — Usuário + Rede",
        ["Ator", "Ação", "Frequência"],
        [
            ["Dona Maria", "Caminhar 15min/dia + alimentação", "Diário"],
            ["Vizinha Joana", "Café da manhã compartilhado", "3x/semana"],
            ["Grupo Igreja", "Acolher no retorno ao artesanato", "Conforme agenda"],
        ],
        LARANJA_ALERTA,
        "PTS é pacto, não prescrição. Ela também tem responsabilidades."
    )

    # ===== SLIDE 24: MOMENTO 4 CONCEITO =====
    add_content_slide(
        prs,
        "④ Reavaliação",
        [
            "PTS é documento VIVO, não estático",
            "Reavaliação PROGRAMADA (não \"se der tempo\")",
            "Avalia: metas alcançadas? ajustes? reclassificar risco?",
            "Critérios: data programada, mudança significativa, meta atingida",
        ],
        ROXO_MOMENTO4,
        "Se metas não estão sendo alcançadas, o que precisa mudar no PTS?"
    )

    # ===== SLIDE 25: CASO ANTÔNIO =====
    add_case_slide(
        prs,
        "Seu Antônio", "68",
        "Transtorno de Ajustamento (F43.21) pós-aposentadoria\nDuração do PTS: 6 meses\n3 reavaliações completadas\n\n✓ 3+ meses sem crise\n✓ Metas alcançadas\n✓ Rede ativada\n✓ Adesão demonstrada",
        "PHQ-9: 12 → 3 (moderada → mínima)",
        "AMARELO → AZUL (alta bem-sucedida)",
        RISCO_AZUL,
        "Antônio é exemplo de ciclo completo. Entrou amarelo, saiu azul."
    )

    # ===== SLIDE 26: CHECKPOINT 3 =====
    add_checkpoint_slide(
        prs,
        "🔄 Reclassificar ou Manter?",
        "Cenário Carla (90 dias):\n• PHQ-9: 18 → 9 (melhora 50%)\n• GAD-7: 14 → 6 (melhora 57%)\n• Metas: 2/3 alcançadas, 1 parcial\n• Boa adesão",
        [
            "Carla deve ser reclassificada de AMARELO para VERDE?",
        ],
        "Discussão: Considerar estabilidade dos ganhos, risco de recaída, suporte social. Reclassificação é evidência de efetividade.",
        "Não há resposta única — depende da avaliação clínica."
    )

    # ===== SLIDE 27: CASO SR. JOÃO =====
    add_case_slide(
        prs,
        "Sr. João", "45",
        "HAS mal controlada\nUso problemático de álcool\nDesempregado há 8 meses\nMora com esposa e 2 filhos adolescentes\n\"Veio só buscar receita de losartana\"\nACS: \"vizinhos dizem que ele bebe todo dia\"",
        "AUDIT: 18 (uso nocivo) | PHQ-9: 11 (moderada)",
        "AMARELO - PTS Obrigatório",
        RISCO_AMARELO,
        "AUDIT 18 = uso nocivo. PHQ-9 11 = depressão moderada (comorbidade comum)."
    )

    # ===== SLIDE 28: INSTRUÇÃO EXERCÍCIO =====
    add_content_slide(
        prs,
        "Tarefa — 10 minutos",
        [
            "☐ Faça o diagnóstico integral (biopsicossocial)",
            "☐ Proponha 2 metas SMART (1 curto, 1 médio prazo)",
            "☐ Defina 3 responsabilidades (equipe + João + família)",
            "",
            ("Lembrem: meta pode ser REDUÇÃO DE DANOS, não abstinência", 18, True, VERDE_SAUDE),
        ],
        AZUL_SUS,
        "Circular entre duplas. Aviso aos 5 e 8 minutos."
    )

    # ===== SLIDE 29: PLENÁRIA =====
    add_content_slide(
        prs,
        "Plenária — 3 duplas (2 min cada)",
        [
            "\"E se não quiser parar de beber?\"",
            ("→ Meta pode ser redução de danos (ex: não beber durante o dia)", 16, False, VERDE_SAUDE),
            "",
            "\"E se a família não colaborar?\"",
            ("→ PTS inclui trabalho com família, mas não depende só dela", 16, False, VERDE_SAUDE),
            "",
            "\"E o desemprego?\"",
            ("→ Meta intersetorial: CRAS, capacitação profissional", 16, False, VERDE_SAUDE),
        ],
        AZUL_SUS,
        "Validar contribuições. Corrigir gentilmente equívocos."
    )

    # ===== SLIDE 30: FLUXOGRAMA =====
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Síntese: Fluxo do PTS"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = AZUL_SUS
    p.alignment = PP_ALIGN.CENTER

    flow_items = [
        ("CASO COMPLEXO", LARANJA_ALERTA),
        ("1. DIAGNÓSTICO INTEGRAL", AZUL_SUS),
        ("2. DEFINIÇÃO DE METAS", VERDE_SAUDE),
        ("3. RESPONSABILIDADES", LARANJA_ALERTA),
        ("4. REAVALIAÇÃO", ROXO_MOMENTO4),
    ]

    y = Inches(1.2)
    for text, color in flow_items:
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), y, Inches(6), Inches(0.7)
        )
        set_shape_fill(box, color)
        box.line.fill.background()

        text_box = slide.shapes.add_textbox(Inches(2), y + Inches(0.15), Inches(6), Inches(0.5))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = BRANCO
        p.alignment = PP_ALIGN.CENTER

        y += Inches(0.9)

        if text != "4. REAVALIAÇÃO":
            arrow = slide.shapes.add_textbox(Inches(4.5), y - Inches(0.25), Inches(1), Inches(0.3))
            tf = arrow.text_frame
            p = tf.paragraphs[0]
            p.text = "↓"
            p.font.size = Pt(24)
            p.font.color.rgb = VERDE_SAUDE
            p.alignment = PP_ALIGN.CENTER

    # ===== SLIDE 31: CHECKLIST =====
    add_content_slide(
        prs,
        "Checklist de Bolso",
        [
            "☐ Caso consumiu múltiplas consultas sem resolução?",
            "☐ Complexidade biopsicossocial evidente?",
            "☐ CuidaSM indica Nível 3 ou 4?",
            "☐ Usuário/família pode participar da construção?",
            "☐ Equipe mínima para dividir responsabilidades?",
            "",
            ("Se SIM para 3 ou mais → considere PTS", 20, True, VERDE_SAUDE),
        ],
        AZUL_SUS,
        "Handout impresso disponível."
    )

    # ===== SLIDE 32: INDICADORES =====
    add_table_slide(
        prs,
        "Indicadores de Qualidade",
        ["Indicador", "Fórmula", "Meta"],
        [
            ["Cobertura PTS CAPS", "(Usuários com PTS / Total) × 100", "100%"],
            ["Participação usuário", "(PTS c/ assinatura / Total) × 100", "≥80%"],
            ["Revisão no prazo", "(Revisados até data / Ativos) × 100", "≥90%"],
            ["Metas alcançadas", "(Alcançadas / Total metas) × 100", "≥70%"],
        ],
        VERDE_SAUDE,
        "Indicadores que vamos acompanhar no município."
    )

    # ===== SLIDE 33: QUIZ =====
    add_quiz_slide(
        prs,
        [
            ("1. PTS é obrigatório para qual classificação mínima?", "Amarelo (Moderado)"),
            ("2. Quantos momentos estruturais tem o PTS?", "4 (Diagnóstico, Metas, Responsabilidades, Reavaliação)"),
            ("3. Metas devem ser definidas por quem?", "Equipe + Usuário (co-construção)"),
            ("4. Qual instrumento usar para estratificar?", "CuidaSM (0-11 pontos, 4 níveis)"),
        ],
        "Se acertaram 3+, o objetivo da aula foi alcançado."
    )

    # ===== SLIDE 34: PRÓXIMOS PASSOS =====
    add_content_slide(
        prs,
        "Próximos Passos",
        [
            ("Esta semana:", 20, True, VERDE_SAUDE),
            "1. Identifique 1 caso \"Dona Maria\" na agenda",
            "2. Aplique os 4 momentos mentalmente",
            "3. Aplique CuidaSM para estratificar",
            "4. Traga dúvidas para o próximo encontro",
            "",
            ("Recursos disponíveis:", 20, True, AZUL_SUS),
            "• Template PTS (F-02) | Checklist de bolso | Escala CuidaSM",
            "• Contato Coordenação SM | Matriciamento mensal",
        ],
        AZUL_SUS,
        "Compromisso concreto. Não deixar abstrato."
    )

    # ===== SLIDE 35: REFERÊNCIAS =====
    add_content_slide(
        prs,
        "Referências",
        [
            ("1. Brasil. MS. Clínica ampliada e PTS. 2ª ed. Brasília: MS; 2007.", 12, False, CINZA_MEDIO),
            ("2. Brasil. MS. CAB 34: Saúde Mental. Brasília: MS; 2013.", 12, False, CINZA_MEDIO),
            ("3. Campos GWS. Apoio matricial. Cad Saúde Pública. 2007;23(2).", 12, False, CINZA_MEDIO),
            ("4. Brasil. MS. Política Nacional de Humanização. Brasília: MS; 2004.", 12, False, CINZA_MEDIO),
            ("5. Brasil. MS. Linha de Cuidado TEA. Brasília: MS; 2025.", 12, False, CINZA_MEDIO),
            ("6. Kroenke K. The PHQ-9. J Gen Intern Med. 2001;16(9):606-13.", 12, False, CINZA_MEDIO),
            ("7. Spitzer RL. The GAD-7. Arch Intern Med. 2006;166(10):1092-7.", 12, False, CINZA_MEDIO),
            ("8. Saunders JB. Development of AUDIT. Addiction. 1993;88(6):791-804.", 12, False, CINZA_MEDIO),
        ],
        AZUL_SUS,
        "Referências em formato Vancouver."
    )

    # ===== SLIDE 36: ENCERRAMENTO =====
    slide = add_title_slide(
        prs,
        "Obrigado!",
        "Vocês na APS são fundamentais.\nA gente está junto nessa."
    )
    bg = slide.shapes[0]
    set_shape_fill(bg, VERDE_SAUDE)

    footer = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1.5))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "5 casos trabalhados: Dona Maria | Lucas | Carla | Seu Antônio | Sr. João"
    p.font.size = Pt(14)
    p.font.color.rgb = BRANCO
    p.alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "\nCoordenação Municipal de Saúde Mental\nExtrema/MG — Janeiro 2026 | v2.0"
    p.font.size = Pt(12)
    p.font.color.rgb = BRANCO
    p.alignment = PP_ALIGN.CENTER

    return prs


def main():
    """Função principal."""
    output_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    output_path = os.path.join(
        output_dir,
        "exports", "aulas", "AULA-01_PTS",
        "AULA-01_PTS_v2.0_2026-01-28.pptx"
    )

    print("Criando apresentação PPTX...")
    prs = create_presentation()

    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)

    print(f"[OK] Apresentacao salva em: {output_path}")
    print(f"     Total de slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
