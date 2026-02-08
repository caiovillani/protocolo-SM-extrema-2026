#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Gerador de Apresentação PowerPoint: Projeto Terapêutico Singular
AULA-01_PTS_v1.0_2026-01-28.pptx

Utiliza python-pptx para criar apresentação seguindo o roteiro detalhado.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Paleta de cores (SUS-inspired)
AZUL_SUS = RGBColor(0x00, 0x5C, 0xA9)
VERDE_SAUDE = RGBColor(0x00, 0xA8, 0x59)
CINZA_ESCURO = RGBColor(0x2D, 0x34, 0x36)
CINZA_CLARO = RGBColor(0xF5, 0xF6, 0xF7)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)
AMARELO_ALERTA = RGBColor(0xFF, 0xC1, 0x07)

def set_slide_background(slide, color):
    """Define cor de fundo do slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs, title, subtitle):
    """Adiciona slide de título."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, BRANCO)

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = AZUL_SUS
    p.alignment = PP_ALIGN.CENTER

    # Subtítulo
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = VERDE_SAUDE
    p.alignment = PP_ALIGN.CENTER

    # Rodapé
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Coordenação Municipal de Saúde Mental — Extrema/MG — Janeiro 2026"
    p.font.size = Pt(14)
    p.font.color.rgb = CINZA_ESCURO
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, bullets, highlight_color=None):
    """Adiciona slide de conteúdo com bullets."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, BRANCO)

    # Barra superior colorida
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = highlight_color or AZUL_SUS
    bar.line.fill.background()

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = AZUL_SUS

    # Bullets
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.5), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(22)
        p.font.color.rgb = CINZA_ESCURO
        p.space_after = Pt(14)

    return slide

def add_case_slide(prs, title, case_text, question=None):
    """Adiciona slide de caso clínico."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, BRANCO)

    # Barra colorida
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = AMARELO_ALERTA
    bar.line.fill.background()

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = AZUL_SUS
    p.alignment = PP_ALIGN.CENTER

    # Caixa do caso
    case_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1.5), Inches(8), Inches(3))
    case_box.fill.solid()
    case_box.fill.fore_color.rgb = CINZA_CLARO
    case_box.line.color.rgb = CINZA_ESCURO

    # Texto do caso
    text_box = slide.shapes.add_textbox(Inches(1.3), Inches(1.8), Inches(7.4), Inches(2.5))
    tf = text_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = case_text
    p.font.size = Pt(22)
    p.font.color.rgb = CINZA_ESCURO
    p.line_spacing = 1.5

    # Pergunta (se houver)
    if question:
        q_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(1))
        tf = q_box.text_frame
        p = tf.paragraphs[0]
        p.text = question
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = VERDE_SAUDE
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_table_slide(prs, title, headers, rows):
    """Adiciona slide com tabela."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, BRANCO)

    # Barra superior
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = AZUL_SUS
    bar.line.fill.background()

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = AZUL_SUS

    # Tabela
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(1.3), Inches(9), Inches(0.6 * num_rows)).table

    # Headers
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = AZUL_SUS
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = BRANCO
        p.alignment = PP_ALIGN.CENTER

    # Rows
    for row_idx, row in enumerate(rows):
        for col_idx, cell_text in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = cell_text
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(14)
            p.font.color.rgb = CINZA_ESCURO
            # Alternating row colors
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CINZA_CLARO

    return slide

def add_diagram_slide(prs, title, steps):
    """Adiciona slide com diagrama de fluxo vertical."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, BRANCO)

    # Barra superior
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = VERDE_SAUDE
    bar.line.fill.background()

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = AZUL_SUS
    p.alignment = PP_ALIGN.CENTER

    # Desenhar boxes para cada step
    y_start = 1.3
    box_height = 0.9
    box_width = 8
    spacing = 0.3

    colors = [AZUL_SUS, VERDE_SAUDE, AZUL_SUS, VERDE_SAUDE]

    for i, (step_title, step_desc) in enumerate(steps):
        y = y_start + i * (box_height + spacing)

        # Box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(y), Inches(box_width), Inches(box_height))
        box.fill.solid()
        box.fill.fore_color.rgb = colors[i % len(colors)]
        box.line.fill.background()

        # Número
        num_box = slide.shapes.add_textbox(Inches(1.2), Inches(y + 0.15), Inches(0.5), Inches(0.6))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = BRANCO

        # Título do step
        step_box = slide.shapes.add_textbox(Inches(1.8), Inches(y + 0.1), Inches(7), Inches(0.4))
        tf = step_box.text_frame
        p = tf.paragraphs[0]
        p.text = step_title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = BRANCO

        # Descrição
        desc_box = slide.shapes.add_textbox(Inches(1.8), Inches(y + 0.5), Inches(7), Inches(0.4))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = step_desc
        p.font.size = Pt(14)
        p.font.color.rgb = BRANCO

        # Seta (exceto último)
        if i < len(steps) - 1:
            arrow_y = y + box_height + 0.05
            arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(4.8), Inches(arrow_y), Inches(0.4), Inches(spacing - 0.1))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = CINZA_ESCURO
            arrow.line.fill.background()

    return slide

def add_checklist_slide(prs, title, items):
    """Adiciona slide de checklist."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, BRANCO)

    # Barra superior
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = VERDE_SAUDE
    bar.line.fill.background()

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = AZUL_SUS

    # Checklist items
    y = 1.5
    for item in items:
        # Checkbox
        checkbox = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(y), Inches(0.35), Inches(0.35))
        checkbox.fill.solid()
        checkbox.fill.fore_color.rgb = BRANCO
        checkbox.line.color.rgb = VERDE_SAUDE

        # Texto
        text_box = slide.shapes.add_textbox(Inches(1.2), Inches(y), Inches(8), Inches(0.5))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = CINZA_ESCURO

        y += 0.7

    return slide

def add_closing_slide(prs, title, next_steps, resources):
    """Adiciona slide de encerramento."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, BRANCO)

    # Barra superior
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = VERDE_SAUDE
    bar.line.fill.background()

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = AZUL_SUS

    # Próximos passos
    steps_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(4.5), Inches(3))
    tf = steps_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Esta semana:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = VERDE_SAUDE

    for step in next_steps:
        p = tf.add_paragraph()
        p.text = f"• {step}"
        p.font.size = Pt(16)
        p.font.color.rgb = CINZA_ESCURO
        p.space_after = Pt(8)

    # Recursos
    res_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.4), Inches(4.3), Inches(3))
    tf = res_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Recursos disponíveis:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = AZUL_SUS

    for res in resources:
        p = tf.add_paragraph()
        p.text = f"• {res}"
        p.font.size = Pt(16)
        p.font.color.rgb = CINZA_ESCURO
        p.space_after = Pt(8)

    # Rodapé
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Obrigado! Dúvidas: Coordenação de Saúde Mental"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = VERDE_SAUDE
    p.alignment = PP_ALIGN.CENTER

    return slide

def create_presentation():
    """Cria a apresentação completa."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # === SLIDE 1: CAPA ===
    add_title_slide(
        prs,
        "Projeto Terapêutico Singular",
        "Uma ferramenta para organizar o cuidado na APS"
    )

    # === SLIDE 2: CASO CLÍNICO - DONA MARIA ===
    add_case_slide(
        prs,
        "Caso Clínico",
        "Dona Maria, 58 anos\nDM2 + HAS + \"depressão\"\n8 consultas no último ano\nSem resolução",
        None
    )

    # === SLIDE 3: PERGUNTA MOBILIZADORA ===
    add_case_slide(
        prs,
        "Reflexão",
        "Dona Maria, 58 anos\nDM2 + HAS + \"depressão\"\n8 consultas no último ano\nSem resolução",
        "Quantas \"Donas Marias\" vocês têm na agenda?"
    )

    # === SLIDE 4: O QUE É PTS? ===
    add_content_slide(
        prs,
        "O que é PTS?",
        [
            "Ferramenta de organização do cuidado",
            "Considera a singularidade de cada caso",
            "Construído COM o usuário, não PARA o usuário",
            "Atribui responsabilidades claras"
        ]
    )

    # === SLIDE 5: PTS vs PLANO GENÉRICO ===
    add_table_slide(
        prs,
        "PTS ≠ Plano de Cuidados Genérico",
        ["Plano Genérico", "PTS"],
        [
            ["Padronizado", "Singularizado"],
            ["Feito pela equipe", "Co-construído com usuário"],
            ["Foco na doença", "Foco na pessoa"],
            ["Atribuições difusas", "Responsabilidades claras"],
            ["Revisão eventual", "Reavaliação programada"]
        ]
    )

    # === SLIDE 6: DE ONDE VEM O PTS? ===
    add_content_slide(
        prs,
        "De onde vem o PTS?",
        [
            "Política Nacional de Humanização (PNH, 2004)",
            "Clínica Ampliada e Compartilhada",
            "Cadernos de Atenção Básica nº 34 (MS, 2013)",
            "Princípio: cuidado centrado na pessoa"
        ],
        VERDE_SAUDE
    )

    # === SLIDE 7: QUANDO ELABORAR PTS? ===
    add_table_slide(
        prs,
        "Quando elaborar PTS?",
        ["Situação", "Indicação"],
        [
            ["Transtorno mental grave e persistente", "Obrigatório"],
            ["Casos AMARELO ou superior", "Obrigatório"],
            ["Compartilhamento APS-AES", "Obrigatório"],
            ["Usuários do CAPS", "Obrigatório"],
            ["Alta complexidade biopsicossocial", "Fortemente recomendado"]
        ]
    )

    # === SLIDE 8: OS 4 MOMENTOS DO PTS ===
    add_diagram_slide(
        prs,
        "Os 4 Momentos do PTS",
        [
            ("DIAGNÓSTICO INTEGRAL", "Bio + Psico + Social + Potencialidades"),
            ("DEFINIÇÃO DE METAS", "SMART — Com o usuário — 30d / 90d / >90d"),
            ("DIVISÃO DE RESPONSABILIDADES", "Equipe + Usuário + Família + Rede"),
            ("REAVALIAÇÃO PROGRAMADA", "Metas alcançadas? Ajustes? Reclassificação?")
        ]
    )

    # === SLIDE 9: DIAGNÓSTICO INTEGRAL - CONCEITO ===
    add_content_slide(
        prs,
        "Momento 1: Diagnóstico Integral",
        [
            "Vai além do CID-10",
            "Mapeia: vulnerabilidades + potencialidades",
            "Avaliação biopsicossocial completa",
            "Inclui: barreiras de acesso, rede de apoio, determinantes sociais"
        ]
    )

    # === SLIDE 10: DIAGNÓSTICO - DONA MARIA ===
    add_table_slide(
        prs,
        "Diagnóstico Integral — Dona Maria",
        ["Dimensão", "Avaliação"],
        [
            ["Biológica", "DM2 descompensada, HAS, insônia crônica"],
            ["Psicológica", "Humor deprimido, luto não elaborado, baixa autoestima"],
            ["Social", "Mora sozinha, filhos distantes, renda insuficiente"],
            ["Potencialidades", "Vínculo com ACS, participava de artesanato"]
        ]
    )

    # === SLIDE 11: DEFINIÇÃO DE METAS - CONCEITO ===
    add_content_slide(
        prs,
        "Momento 2: Definição de Metas",
        [
            "Metas SMART: Específicas, Mensuráveis, Alcançáveis, Relevantes, Temporais",
            "Negociadas COM o usuário — não impostas",
            "Horizonte: curto (30d), médio (90d), longo prazo (>90d)",
            "Foco em funcionalidade e qualidade de vida"
        ],
        VERDE_SAUDE
    )

    # === SLIDE 12: METAS - DONA MARIA ===
    add_table_slide(
        prs,
        "Metas SMART — Dona Maria",
        ["Prazo", "Meta", "Indicador"],
        [
            ["30 dias", "Retomar sono de 6h/noite", "Autorrelato"],
            ["30 dias", "Rotina de café com vizinha", "Confirmação ACS"],
            ["90 dias", "Retornar ao artesanato", "Frequência ≥2x/mês"],
            ["90 dias", "HbA1c < 8%", "Exame laboratorial"]
        ]
    )

    # === SLIDE 13: DIVISÃO DE RESPONSABILIDADES - CONCEITO ===
    add_content_slide(
        prs,
        "Momento 3: Divisão de Responsabilidades",
        [
            "QUEM faz O QUÊ, QUANDO",
            "Inclui: equipe + usuário + família + rede",
            "Define profissional de referência (gestor do caso)",
            "Evita: 'todo mundo cuida = ninguém cuida'"
        ]
    )

    # === SLIDE 14: RESPONSABILIDADES - DONA MARIA ===
    add_table_slide(
        prs,
        "Responsabilidades — Dona Maria",
        ["Ator", "Ação", "Frequência"],
        [
            ["Médico(a) de Família", "Revisão medicamentosa + humor", "Mensal"],
            ["Enfermeiro(a)", "Monitoramento DM/HAS + escuta", "Quinzenal"],
            ["ACS", "Visita + monitoramento sono", "Semanal"],
            ["Psicólogo (eMulti)", "Acompanhamento do luto", "Quinzenal"],
            ["Dona Maria", "Caminhar 15min + registro alimentar", "Diário"],
            ["Vizinha", "Café da manhã compartilhado", "3x/semana"]
        ]
    )

    # === SLIDE 15: PROFISSIONAL DE REFERÊNCIA ===
    add_content_slide(
        prs,
        "Profissional de Referência",
        [
            "É quem mantém o vínculo principal",
            "Coordena o PTS (não faz tudo sozinho)",
            "Geralmente: médico ou enfermeiro da ESF",
            "Vocês são os melhores candidatos — conhecem o território"
        ],
        VERDE_SAUDE
    )

    # === SLIDE 16: REAVALIAÇÃO - CONCEITO ===
    add_content_slide(
        prs,
        "Momento 4: Reavaliação",
        [
            "PTS é documento VIVO, não estático",
            "Reavaliação programada (não apenas 'se der tempo')",
            "Avalia: metas alcançadas? ajustes? reclassificação?",
            "Pode resultar em: manutenção, intensificação ou alta"
        ]
    )

    # === SLIDE 17: EXERCÍCIO GUIADO - SR. JOÃO ===
    add_case_slide(
        prs,
        "Exercício Guiado — Caso Sr. João",
        "Sr. João, 45 anos\nHAS + Uso problemático de álcool\nDesempregado há 8 meses\nMora com esposa e 2 filhos adolescentes\nÚltima consulta: \"veio só buscar receita de losartana\"\nACS relata: \"Vizinhos comentam que ele bebe todo dia\"",
        None
    )

    # === SLIDE 18: TAREFA EM DUPLAS ===
    add_content_slide(
        prs,
        "Tarefa (em duplas, 8 minutos)",
        [
            "1. Faça o diagnóstico integral (biopsicossocial)",
            "2. Proponha 2 metas SMART (1 curto, 1 médio prazo)",
            "3. Defina 3 responsabilidades (equipe + João + família)"
        ],
        AMARELO_ALERTA
    )

    # === SLIDE 19: SÍNTESE VISUAL ===
    add_diagram_slide(
        prs,
        "Síntese: Fluxo PTS",
        [
            ("CASO COMPLEXO", "Múltiplas consultas sem resolução"),
            ("DIAGNÓSTICO INTEGRAL", "Bio + Psico + Social + Potencialidades"),
            ("METAS SMART", "Negociadas com usuário — 30/90/>90 dias"),
            ("RESPONSABILIDADES", "Quem faz o quê, quando"),
            ("REAVALIAÇÃO", "Metas alcançadas? Ajustes? Alta?")
        ]
    )

    # === SLIDE 20: CHECKLIST DE BOLSO ===
    add_checklist_slide(
        prs,
        "Checklist: Quando elaborar PTS?",
        [
            "Este caso já consumiu múltiplas consultas sem resolução?",
            "Há complexidade biopsicossocial evidente?",
            "O usuário/família pode participar da construção?",
            "Tenho equipe mínima para dividir responsabilidades?",
            "",
            "Se SIM para 3 ou mais → considere PTS"
        ]
    )

    # === SLIDE 21: PRÓXIMOS PASSOS ===
    add_closing_slide(
        prs,
        "Próximos Passos",
        [
            "Identifique 1 caso \"Dona Maria\" na sua agenda",
            "Aplique os 4 momentos mentalmente",
            "Traga dúvidas para o próximo encontro"
        ],
        [
            "Template PTS simplificado (F-02)",
            "Contato da Coordenação de SM",
            "Matriciamento mensal com eMulti/NIRSM"
        ]
    )

    return prs

def main():
    """Função principal."""
    output_dir = Path(__file__).parent.parent / "exports" / "aulas" / "AULA-01_PTS"
    output_dir.mkdir(parents=True, exist_ok=True)

    output_file = output_dir / "AULA-01_PTS_v1.0_2026-01-28.pptx"

    print("Gerando apresentação PowerPoint...")
    prs = create_presentation()
    prs.save(str(output_file))
    print(f"Apresentação salva em: {output_file}")
    print(f"Total de slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
